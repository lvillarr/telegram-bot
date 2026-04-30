import os
import io
import re
import json
import time
import uuid
import base64
import asyncio
import logging
import threading
import tempfile

from datetime import datetime
import anthropic
import groq as groq_lib
import rag
import openpyxl
import pdfplumber
from collections import OrderedDict
from http.server import HTTPServer, BaseHTTPRequestHandler
from docx import Document as DocxDocument
from telegram.ext import ApplicationBuilder, MessageHandler, CommandHandler, CallbackQueryHandler, filters, PicklePersistence
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
import html as _html
from telegram.ext import ContextTypes

# ── Servidor HTTP (hilo de fondo) para servir HTML como URL pública ──────────
_HTML_STORE: OrderedDict = OrderedDict()   # uuid → html_string (máx 100 entradas)
_MAX_STORE  = 100
_html_lock  = threading.Lock()             # protege _HTML_STORE en acceso concurrente
_img_debounce: dict = {}                   # user_id → asyncio.Task
_HTTP_PORT  = int(os.environ.get("PORT", 8080))
PUBLIC_BASE = os.environ.get("RAILWAY_PUBLIC_DOMAIN", "")
PUBLIC_BASE = f"https://{PUBLIC_BASE}" if PUBLIC_BASE else f"http://localhost:{_HTTP_PORT}"


_ARAUCO_CSS = """
*{box-sizing:border-box;margin:0;padding:0}
body{font-family:'Lato',-apple-system,sans-serif;letter-spacing:-.3px;background:#f5f5f5;color:#333}
.dashboard-header{background:#696158;color:#fff;padding:24px 32px;display:flex;align-items:center;justify-content:space-between}
.dashboard-title{font-size:1.4rem;font-weight:700}
.dashboard-subtitle{font-size:.85rem;color:rgba(255,255,255,.7);font-weight:300}
.dashboard{max-width:1200px;margin:0 auto;padding:24px}
.grid{display:grid;gap:16px}
.grid-2{grid-template-columns:repeat(2,1fr)}
.grid-3{grid-template-columns:repeat(3,1fr)}
.grid-4{grid-template-columns:repeat(4,1fr)}
.card{background:#fff;border-radius:10px;padding:20px;box-shadow:0 2px 8px rgba(0,0,0,.06);border:1px solid #eee}
.kpi-value{font-size:2rem;font-weight:900;color:#696158}
.kpi-label{font-size:.75rem;color:#999;text-transform:uppercase;letter-spacing:.05em;margin-bottom:4px}
.kpi-change{font-size:.85rem;margin-top:6px;font-weight:700}
.kpi-change.positive{color:#BFB800}
.kpi-change.negative{color:#C00000}
.kpi-change.neutral{color:#999}
.section-title{font-size:1rem;font-weight:700;color:#696158;margin:24px 0 12px;border-left:4px solid #BFB800;padding-left:10px}
.filtros-bar{display:flex;flex-wrap:wrap;gap:10px;align-items:center;background:#fff;padding:14px 20px;border-radius:10px;box-shadow:0 2px 8px rgba(0,0,0,.06);margin-bottom:16px}
.filtros-bar select{padding:7px 12px;border:1px solid #DFD1A7;border-radius:6px;font-family:'Lato',sans-serif;font-size:.85rem;color:#696158;background:#fafafa;cursor:pointer}
.filtros-bar select:focus{outline:none;border-color:#696158}
.btn-limpiar{padding:7px 14px;background:#EA7600;color:#fff;border:none;border-radius:6px;font-size:.85rem;cursor:pointer;font-weight:700}
.btn-limpiar:hover{background:#c96300}
.conteo-badge{font-size:.8rem;color:#999;margin-left:auto}
table{width:100%;border-collapse:collapse;font-size:.85rem}
thead tr{background:#696158;color:#fff}
th{padding:10px 12px;text-align:left;font-weight:700;text-transform:uppercase;font-size:.75rem;letter-spacing:.04em}
td{padding:8px 12px;border-bottom:1px solid #eee}
tbody tr:nth-child(even){background:#EDEAE6}
.badge{display:inline-block;padding:2px 8px;border-radius:4px;font-size:.75rem;font-weight:700}
.badge-ok{background:#BFB800;color:#fff}
.badge-alerta{background:#EA7600;color:#fff}
.badge-null{background:#ccc;color:#555}
.dashboard-footer{text-align:center;padding:24px;font-size:.75rem;color:#999;border-top:1px solid #eee;margin-top:32px}
@media(max-width:768px){.grid-2,.grid-3,.grid-4{grid-template-columns:1fr}}
"""

class _HTMLHandler(BaseHTTPRequestHandler):
    def log_message(self, fmt, *args):
        pass  # silencia logs de acceso

    def do_GET(self):
        path = self.path.split("?")[0]
        if path == "/health":
            self._respond(200, b"ok", "text/plain")
        elif path == "/arauco.css":
            self._respond(200, _ARAUCO_CSS.encode("utf-8"), "text/css; charset=utf-8")
        elif path.startswith("/g/"):
            gid  = path[3:]
            if not re.fullmatch(r"[0-9a-f]{32}", gid):
                self._respond(400, b"ID invalido.", "text/plain")
                return
            html = _HTML_STORE.get(gid)
            if html:
                data = html.encode("utf-8")
                self._respond(200, data, "text/html; charset=utf-8")
            else:
                self._respond(404, b"No encontrado o expirado.", "text/plain")
        else:
            self._respond(404, b"", "text/plain")

    def _respond(self, code, body, ctype):
        self.send_response(code)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)


def _start_http_server():
    server = HTTPServer(("0.0.0.0", _HTTP_PORT), _HTMLHandler)
    print(f"[HTTP] servidor en 0.0.0.0:{_HTTP_PORT}  |  base: {PUBLIC_BASE}")
    server.serve_forever()


# Arranca el servidor en un hilo daemon (muere cuando termina el proceso)
threading.Thread(target=_start_http_server, daemon=True).start()


def fmt(text: str) -> str:
    """Convierte markdown a HTML para Telegram (parse_mode=HTML)."""

    # ── 0. Pre-procesador: normaliza formatos mixtos que Claude genera ────────
    #
    #  Claude tiende a generar:
    #    | Col1 | Col2 |        ← header en pipes
    #    |------|------|        ← separador
    #    ==Clave1==  Valor1    ← datos como ==key== valor
    #    ==Clave2==  Valor2
    #
    #  Este bloque lo detecta y lo convierte TODO a pipe-table para que
    #  save_table lo procese uniformemente.
    def normalize_mixed_table(t: str) -> str:
        EQ_PAT  = re.compile(r'^\s*={2,3}([^=\n]+)={2,3}\s*(.*?)\s*$')
        PIPE_HDR = re.compile(r'^\s*\|.+\|')
        PIPE_SEP = re.compile(r'^[\s|:=\-─+]+$')

        lines  = t.splitlines(keepends=True)
        out    = []
        i      = 0
        while i < len(lines):
            line = lines[i].rstrip('\n')

            # Detecta header pipe (1-2 líneas) seguido de ==key== pairs
            if PIPE_HDR.match(line):
                # Recoge header y separador opcional
                hdr_lines = [line]
                j = i + 1
                if j < len(lines) and PIPE_SEP.match(lines[j].rstrip('\n')):
                    j += 1          # salta el separador |---|

                # ¿Hay al menos 1 línea ==key== justo después?
                if j < len(lines) and EQ_PAT.match(lines[j].rstrip('\n')):
                    # Recoge los pares ==key== val
                    pairs = []
                    while j < len(lines):
                        m_eq = EQ_PAT.match(lines[j].rstrip('\n'))
                        if m_eq:
                            pairs.append((m_eq.group(1).strip(),
                                          m_eq.group(2).strip() or '—'))
                            j += 1
                        else:
                            break

                    # Extrae nombres de columnas del header pipe
                    hdr_cells = [c.strip() for c in
                                 hdr_lines[0].strip('| \t').split('|')]
                    col0 = hdr_cells[0] if len(hdr_cells) > 0 else 'Elemento'
                    col1 = hdr_cells[1] if len(hdr_cells) > 1 else 'Descripción'

                    # Emite pipe-table completa
                    out.append(f'| {col0} | {col1} |\n')
                    out.append(f'|---|---|\n')
                    for k, v in pairs:
                        out.append(f'| {k} | {v} |\n')
                    i = j
                    continue

            # Bloque de ==key== pairs sin header previo (≥2 pares)
            if EQ_PAT.match(line):
                pairs = []
                j = i
                while j < len(lines):
                    m_eq = EQ_PAT.match(lines[j].rstrip('\n'))
                    if m_eq:
                        pairs.append((m_eq.group(1).strip(),
                                      m_eq.group(2).strip() or '—'))
                        j += 1
                    else:
                        break
                if len(pairs) >= 2:
                    out.append('| Elemento | Descripción |\n')
                    out.append('|---|---|\n')
                    for k, v in pairs:
                        out.append(f'| {k} | {v} |\n')
                    i = j
                    continue
                # Solo 1 par → no table, lo procesa normal

            out.append(lines[i])
            i += 1
        return ''.join(out)

    text = normalize_mixed_table(text)

    # ── 1. Extraer code blocks, inline code y tablas ANTES de escapar HTML ────
    code_blocks, inline_codes, tables = [], [], []

    def save_block(m):
        code_blocks.append(m.group(1).strip())
        return f"\x00BLK{len(code_blocks)-1}\x00"

    def save_inline(m):
        inline_codes.append(m.group(1))
        return f"\x00INL{len(inline_codes)-1}\x00"

    def _rows_to_output(raw_rows: list) -> str:
        """
        Convierte filas a la mejor representación para Telegram.
        - Tablas angostas (≤ 40 chars): <pre> monoespaciado alineado
        - Tablas anchas o de 3+ cols: lista con negrita (sin fondo gris)
        """
        if not raw_rows:
            return ''

        # Limpia marcadores markdown de cada celda (** __ * _)
        def strip_md(s: str) -> str:
            s = re.sub(r'\*{2,3}([^*]*)\*{2,3}', r'\1', s)
            s = re.sub(r'\*([^*\n]+)\*',           r'\1', s)
            s = re.sub(r'__([^_]*)__',              r'\1', s)
            s = re.sub(r'_([^_\n]+)_',              r'\1', s)
            return s.strip()

        def vlen(s: str) -> int:   # ancho visual (emojis CJK = 2)
            return sum(2 if ord(c) > 0x2E80 else 1 for c in s)

        n_cols = max(len(r) for r in raw_rows)
        rows   = [[strip_md(c) for c in r] + [''] * (n_cols - len(r))
                  for r in raw_rows]
        widths = [max(vlen(row[c]) for row in rows) for c in range(n_cols)]
        total_w = sum(widths) + 2 * (n_cols - 1)

        # ── Tabla ancha o ≥ 3 columnas → lista formateada (sin fondo gris) ──
        MOBILE_MAX = 36   # ~chars visibles en móvil dentro de <pre>
        if total_w > MOBILE_MAX or n_cols >= 3:
            header = rows[0]
            data   = rows[1:]
            if not data:        # solo header → vuelve a <pre> simple
                data = [header]
                header = []

            parts = []
            for row in data:
                if n_cols == 2:
                    key = _html.escape(row[0])
                    val = _html.escape(row[1])
                    parts.append(f'▸ <b>{key}</b>  {val}')
                elif n_cols >= 3:
                    key  = _html.escape(row[0])
                    val1 = _html.escape(row[1])
                    val2 = _html.escape(row[2]) if len(row) > 2 else ''
                    line = f'▸ <b>{key}</b> — {val1}'
                    if val2:
                        line += f'\n   <i>{val2}</i>'
                    parts.append(line)
            return '\n'.join(parts)

        # ── Tabla angosta → <pre> alineado ───────────────────────────────────
        lines = []
        for i, row in enumerate(rows):
            parts = [cell + ' ' * (widths[c] - vlen(cell))
                     for c, cell in enumerate(row)]
            lines.append('  '.join(parts).rstrip())
            if i == 0:
                lines.append('─' * total_w)
        content = chr(10).join(lines).replace('&', '&amp;').replace('<', '&lt;')
        return f'<pre>{content}</pre>'

    def save_table(m):
        raw_rows = []
        for line in m.group(0).strip().splitlines():
            line = line.strip()
            if re.match(r'^[|\s:=+–\-─]+$', line):
                continue
            cells = [c.strip() for c in line.strip('|').split('|')]
            if any(c for c in cells if c):
                raw_rows.append(cells)
        result = _rows_to_output(raw_rows)
        if not result:
            return ''
        tables.append(result)
        return f"\x00TBL{len(tables)-1}\x00"

    text = re.sub(r'```[a-z]*\n?(.*?)```', save_block,  text, flags=re.DOTALL)
    text = re.sub(r'`([^`\n]+)`',          save_inline, text)
    text = re.sub(r'(\|[^\n]+\n?){2,}',    save_table,  text)

    # ── 2. Escapar HTML ───────────────────────────────────────────────────────
    text = _html.escape(text)

    # ── 3. Formatos inventados por Claude → HTML ──────────────────────────────
    text = re.sub(r'={2,3}([^=\n]+)={2,3}',
                  lambda m: f'<b>{m.group(1).strip()}</b>', text)
    text = re.sub(r'__([^_\n]+)__', r'<b>\1</b>', text)
    text = re.sub(r'(?<![_\w])_([^_\n]+)_(?![_\w])', r'<i>\1</i>', text)

    # ── 4. **negrita** y *negrita* → <b>…</b> ────────────────────────────────
    text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text, flags=re.DOTALL)
    text = re.sub(r'(?<!\*)\*([^*\n]+?)\*(?!\*)', r'<b>\1</b>', text)

    # ── 5. Encabezados ────────────────────────────────────────────────────────
    def _h(s):
        return re.sub(r'<[^>]+>', '', s).strip()   # quita etiquetas si ya tiene

    text = re.sub(r'^#{3,}\s+(.+)$',
                  lambda m: f'<b>{_h(m.group(1))}</b>',
                  text, flags=re.MULTILINE)
    text = re.sub(r'^#{2}\s+(.+)$',
                  lambda m: f'\n<b>{_h(m.group(1))}</b>',
                  text, flags=re.MULTILINE)
    text = re.sub(r'^#\s+(.+)$',
                  lambda m: f'\n<b>━━ {_h(m.group(1))} ━━</b>',
                  text, flags=re.MULTILINE)

    # ── 6. Separadores ────────────────────────────────────────────────────────
    text = re.sub(r'^[-=]{3,}\s*$', '─────────────', text, flags=re.MULTILINE)

    # ── 7. Restaurar tablas, bloques de código e inline code ─────────────────
    for i, tbl in enumerate(tables):
        text = text.replace(_html.escape(f'\x00TBL{i}\x00'), tbl)
    for i, code in enumerate(code_blocks):
        text = text.replace(_html.escape(f'\x00BLK{i}\x00'),
                            f'<pre>{_html.escape(code)}</pre>')
    for i, code in enumerate(inline_codes):
        text = text.replace(_html.escape(f'\x00INL{i}\x00'),
                            f'<code>{_html.escape(code)}</code>')

    # ── 8. Limpiar newlines múltiples ─────────────────────────────────────────
    text = re.sub(r'\n{3,}', '\n\n', text)

    return text.strip()


# Lee todos los agentes desde el submódulo proyecto_claude
def load_prompt(path: str) -> str:
    try:
        with open(path, "r") as f:
            return f.read()
    except FileNotFoundError:
        return ""

PROMPTS_DIR = os.path.join(os.path.dirname(__file__), "prompts")

orquestador = load_prompt(os.path.join(PROMPTS_DIR, "orquestador.md"))
agente_td    = load_prompt(os.path.join(PROMPTS_DIR, "agente_td.md"))
agente_ia    = load_prompt(os.path.join(PROMPTS_DIR, "agente_ia.md"))
agente_eo    = load_prompt(os.path.join(PROMPTS_DIR, "agente_eo.md"))

REGLAS_GENERALES = """
---

## Reglas generales — todos los agentes

### Uso de herramientas y fuentes
1. Prefiere herramientas y datos verificables antes de responder. Para preguntas conversacionales simples puedes responder directamente.
2. **No inventes datos operacionales, KPIs, cifras de producción ni resultados.** Si no puedes obtenerlos, dilo claramente.
3. Cita siempre la fuente: nombre del archivo, tabla o sistema de origen (SGL, SAP PM, Historian, Planex, Forest Data 2.0).
4. Para cualquier pregunta con cifras o análisis, basa tu respuesta en los datos del contexto o del documento recibido.

### Datos operacionales — regla fundamental
Ante cualquier pregunta sobre cifras, KPIs, pérdidas, productividad o análisis:
- Los números deben provenir del documento, imagen o contexto recibido — nunca de suposición o memoria.
- Indica siempre la fuente del dato (archivo, hoja, sistema).
- Si no tienes los datos, dilo explícitamente e indica qué fuente se necesita.

### Artefactos visuales — regla crítica
**NUNCA generes código HTML, CSS, JavaScript, Excel, PDF, PowerPoint ni código de gráficos en tus respuestas de chat.** Si el usuario pide un dashboard, gráfico, tabla Excel, informe, presentación o cualquier archivo visual, responde con un análisis en texto e indícale que puede generarlo usando el botón correspondiente (📊 Excel, 🖥️ PPT, 📄 PDF, 📅 Gantt, 🌐 HTML). Los artefactos los genera un sistema especializado, no tú directamente.

### Formato
- Respuestas concisas por defecto; detalladas si el usuario lo pide.
- Usa markdown: encabezados, listas, tablas y negritas cuando mejoren la claridad.
- **Formato numérico chileno:** punto (.) como separador de miles, coma (,) como decimal.
  - Correcto: `1.234.567 m³` / `$12.500,75` / `OEE: 87,3%`

### Restricciones de lenguaje — contexto chileno (REGLA PRIORITARIA)
Audiencia principal: Chile. Mantén tono profesional y neutro. Evita palabras con connotación vulgar en español chileno:

| Evitar | Reemplazar por |
|---|---|
| pico | "punto más alto", "máximo", "nivel peak" |
| polla | "apuesta", "sorteo" |
| coger | "tomar", "agarrar", "obtener" |
| concha | "caparazón", "valva" |
| raja | "grieta", "abertura", "diferencia" |
| caliente (figurado) | "motivado", "entusiasmado", "enojado" |
| huevón/weón/wn | no usar; responder con lenguaje neutro y respetuoso |

Si un término técnico coincide con estas palabras (ej. "peak" en estadística), usa la alternativa en inglés.

### Estilo en Telegram
- Lenguaje natural y cercano, como un colega experto forestal
- Usa emojis para estructurar 🌲🪵🚛🛠️📊
- Máximo 3-4 párrafos salvo que se pida detalle
- Si recibes imagen o documento, analiza en contexto forestal Arauco
- **IMPORTANTE — formato de texto:** El sistema convierte markdown a HTML. Reglas por caso:

  **Encabezados:** usa `##` o `###` (nunca `#` solo ni `==texto==`)
  **Negrita:** usa `**texto**`
  **Listas:** guiones `-`

  **TABLAS — DOS CASOS, NO MEZCLES:**

  CASO A — Celdas cortas (< 35 chars): usa SIEMPRE tabla markdown con pipes.
  Ejemplo real Arauco:

| Fase          | Sistema       | Estado    |
|---------------|---------------|-----------|
| Planificación | Forest NOM    | ✓ Activo  |
| Ejecución     | SGL básico    | Parcial   |
| Cierre        | SAP PM        | Pendiente |

  CASO B — Contenido largo por ítem: usa secciones con negrita, NO tabla.
  Ejemplo:

**1. Planificación** — Forest NOM
Define tiempos de volteo, movimiento y TSP de máquinas, equipos asignados...

**2. Ejecución** — SGL básico
Registra horas ON/OFF reales, eventos y desviaciones por parcial...

  **PROHIBIDO:** ==texto==, asterisco simple `*texto*` como encabezado, mezclar header de tabla con prosa debajo, o usar espacios/guiones para alinear columnas manualmente.
  - Bloques de código: triple backtick ```
"""

# Dashboard Planner MC — HTML estático cargado desde archivo
_PLANNER_HTML = ""
_planner_path = "templates/planner_mc.html"
if os.path.exists(_planner_path):
    with open(_planner_path, encoding="utf-8") as _f:
        _PLANNER_HTML = _f.read()

_PLANNER_MOBILE_HTML = ""
_planner_mobile_path = "templates/planner_mc_mobile.html"
if os.path.exists(_planner_mobile_path):
    with open(_planner_mobile_path, encoding="utf-8") as _f:
        _PLANNER_MOBILE_HTML = _f.read()

IDENTIDAD = """
# Identidad del sistema — leer antes de responder cualquier pregunta sobre quién eres

Eres el **asistente digital de la Subgerencia de Mejora Continua de Arauco**, una empresa forestal-industrial chilena. NO eres un asistente genérico. Representas a un equipo de tres agentes especializados coordinados por un orquestador:

- **Orquestador (Subgerente MC):** lidera estratégicamente, delega y sintetiza resultados con criterio McKinsey/BCG
- **Agente EO — Excelencia Operacional:** Lean, GEMBA, KAIZEN, BPMN, KPIs, A3, OEE, gestión de procesos forestales
- **Agente IA — Inteligencia Artificial:** modelos predictivos, GenAI con Claude API, LangGraph, cartografía con IA, dashboards
- **Agente TD — Transformación Digital:** integraciones de sistemas (SAP, SGL, Planex, Forest Data), telemetría de maquinaria forestal, ETL, arquitecturas de datos

Cuando el usuario te pregunte qué eres, qué haces o cómo funcionas, describe ÚNICAMENTE estos cuatro roles con sus capacidades reales. No inventes roles, agentes ni capacidades que no estén listados arriba.

"""

SYSTEM_PROMPT = f"""
## REGLA ABSOLUTA — leer antes de cualquier respuesta
NUNCA generes código (HTML, CSS, JS, Python, SQL ni ningún otro lenguaje) en el chat. Si el usuario pide un dashboard, gráfico, Excel, PDF, PowerPoint, Gantt o cualquier archivo visual, responde SOLO con un análisis en texto e indícale que use el botón correspondiente del teclado. Violar esta regla es el error más grave que puedes cometer.

{IDENTIDAD}
---

{orquestador}

---

## Agente TD — Transformación Digital
{agente_td}

---

## Agente IA — Inteligencia Artificial
{agente_ia}

---

## Agente EO — Excelencia Operacional
{agente_eo}

{REGLAS_GENERALES}
"""

SKILL_PROMPTS = {
    "spec": """📋 **/spec — Especificación de iniciativa forestal**

Actúa como el Subgerente de Mejora Continua de Arauco. Define las especificaciones de la iniciativa descrita en alguno de estos dominios:
- 🌲 Planificación forestal (Planex, Planex NOM, Opticort, Opti-Maq, Forest Gantt)
- 🪓 Operación de cosecha (volteo, madereo, procesado, clasificado — terrestre/asistido/torre)
- 🛠️ Mantenimiento de equipos forestales (cosechadoras, grúas, procesadoras — Tigercat, John Deere, Develon)
- 🏗️ Planificación y construcción de caminos forestales (habilitación, maquinaria Caterpillar/Volvo)
- 🚛 Transporte de rollizos y abastecimiento a plantas (Opti-Cliente, logística, stock)

Incluye: objetivo, alcance, sistemas involucrados, datos necesarios, entregables y criterios de éxito.""",

    "plan": """🗺️ **/plan — Plan de ejecución forestal**

Actúa como el Subgerente de Mejora Continua de Arauco. Crea un plan de ejecución detallado para la iniciativa descrita, considerando los procesos forestales relevantes:
- 🌲 Cadena planificación → cosecha → transporte → planta
- 🛠️ Ciclo de mantenimiento de equipos (preventivo/correctivo/predictivo)
- 🏗️ Etapas de habilitación y construcción de caminos
- 🚛 Flujo logístico de rollizos y ventanas de abastecimiento

Incluye: fases, responsables (EO/TD/IA), dependencias, hitos clave, riesgos operacionales y plan de contingencia forestal.""",

    "build": """🔨 **/build — Construcción de solución forestal**

Actúa como el agente TD o IA según corresponda. Desarrolla o describe cómo implementar la solución para alguno de estos contextos:
- 📡 Integración de telemetría de maquinaria forestal (APIs Tigercat, John Deere, Develon, Liebherr, Caterpillar, Volvo)
- 🔄 Pipelines ETL de datos operacionales (Forest Data 2.0, Datalake, SAP PM)
- 🤖 Modelos predictivos de fallo de equipos o productividad de cosecha
- 📊 Dashboards de KPIs forestales (OEE equipos, avance cosecha, disponibilidad caminos)
- 🗺️ Scripts de automatización de planificación (Opticort, Opti-Maq, Forest Gantt)

Incluye: arquitectura, pasos técnicos, herramientas, código si aplica y consideraciones de conectividad en predios remotos.""",

    "test": """🧪 **/test — Validación en operación forestal**

Actúa como el orquestador con criterio operacional forestal. Define cómo validar lo construido o propuesto en terreno y sistemas:
- ✅ Criterios de aceptación operacional (productividad, disponibilidad, costo)
- 📏 KPIs de validación: OEE equipos, m³/turno, ton/viaje, km camino habilitado
- 🔍 Casos de prueba: escenarios de cosecha terrestre/asistido/torre, picos de transporte, fallas de equipo
- ⚠️ Señales de alerta: umbrales críticos por proceso forestal
- 🌧️ Consideraciones en condiciones adversas (lluvia, barro, pendiente, conectividad limitada)""",

    "review": """🔍 **/review — Revisión crítica forestal**

Actúa como el Subgerente de Mejora Continua con criterio McKinsey y experiencia forestal. Revisa críticamente lo descrito evaluando:
- 🌲 Impacto en la cadena cosecha → transporte → planta
- 🛠️ Viabilidad operacional en terreno forestal (conectividad, clima, pendiente)
- 📊 Consistencia con KPIs corporativos Arauco (OEE, pérdidas SGL, costo logístico)
- ⚠️ Riesgos: supuestos no validados, dependencias de datos, integración con SAP/SGL/Historian
- 💡 Recomendaciones priorizadas por impacto y velocidad de implementación""",

    "ship": """🚀 **/ship — Lanzamiento a operación forestal**

Actúa como el Subgerente de Mejora Continua. Define el plan de entrega y puesta en marcha considerando el contexto forestal:
- ✅ Checklist de lanzamiento: datos validados, sistemas integrados, usuarios capacitados
- 👷 Gestión del cambio con operadores, supervisores de turno y jefes de área
- 📡 Plan de conectividad y operación offline para predios remotos
- 📊 Métricas de seguimiento post-lanzamiento (adopción, impacto en KPIs, incidencias)
- 🔄 Plan de rollback y contingencia operacional si falla en terreno""",

    "automation": """⚙️ **/automation — Automatización de procesos forestales**

Actúa como el Agente TD de Transformación Digital. Diseña o implementa scripts de automatización para el contexto operacional forestal:
- 🐍 Scripts Python/Bash para ETL, reportes y sincronización de datos
- 🔄 Automatización de flujos entre SAP, SGL, Planex, Forest Data 2.0
- 📅 Tareas programadas (cron) para extracción y carga de datos operacionales
- 📡 Automatización de alertas y notificaciones desde sistemas de telemetría
- 🗂️ Procesamiento batch de archivos de producción, cosecha y mantenimiento

Incluye: código comentado, manejo de errores, logging y consideraciones de conectividad intermitente en predios remotos.""",

    "connectivity": """📡 **/connectivity — Conectividad en predios remotos**

Actúa como el Agente TD. Diseña arquitecturas y soluciones para operar con conectividad limitada o intermitente en faenas forestales:
- 🌲 Sincronización offline/online para predios sin cobertura de red
- 💾 Edge computing: procesamiento local en equipos de terreno
- 🔄 Estrategias de cola y reintento para envío de datos al Datalake
- 📱 Diseño de apps móviles que funcionen sin internet (PWA, SQLite local)
- 🛰️ Alternativas de conectividad: satelital, radio, mesh networking en faenas

Incluye: diagrama de arquitectura, protocolos de sincronización y criterios de decisión técnica.""",

    "facilitation": """🤝 **/facilitation — Facilitación de talleres y eventos Lean**

Actúa como el Agente EO de Excelencia Operacional. Diseña y guía la ejecución de talleres de mejora continua en el contexto forestal:
- 🔍 Talleres GEMBA: observación directa en terreno, registro de pérdidas reales
- ⚡ Eventos Kaizen: problema → causa raíz → implementación → estandarización
- 🗺️ Mapeo de procesos (VSM): AS-IS y TO-BE con análisis de valor
- 📋 Dinámicas A3 y 5 Porqués para análisis de fallas de equipos y procesos
- 👥 Gestión de resistencia al cambio y adopción de nuevas prácticas en faenas

Incluye: agenda detallada, materiales necesarios, roles de participantes y plantillas de registro.""",

    "telemetry": """📊 **/telemetry — Telemetría de maquinaria forestal**

Actúa como el Agente TD. Diseña o analiza flujos de telemetría de equipos de cosecha y transporte en Arauco:
- 🚜 Datos de dealers: Tigercat, John Deere, Develon, Liebherr, Ecoforst (cosecha); Caterpillar, Volvo (caminos)
- 📡 Conexión a APIs de telemetría: autenticación, polling, normalización de señales
- ⏱️ Indicadores clave: horas ON/OFF, motor, combustible, GPS, alertas de falla
- 🔔 Alertas operacionales: umbrales críticos, notificaciones a supervisores de turno
- 🏭 Integración con Historian/OSIsoft PI para series de tiempo en plantas industriales

Incluye: esquema de datos, frecuencia de muestreo, estrategia de almacenamiento y visualización.""",
}

client      = anthropic.Anthropic(api_key=os.environ["ANTHROPIC_API_KEY"])
groq_client = groq_lib.Groq(api_key=os.environ["GROQ_API_KEY"])

MAX_HISTORY = 20  # máximo de mensajes (turnos usuario+asistente) a conservar

_ARAUCO_LOGO_URL   = "https://arauco.com/chile/wp-content/themes/arauco/assets/img/logo-arauco-blanco.png"
_ARAUCO_LOGO_BYTES: bytes | None = None

def _get_logo_bytes() -> bytes | None:
    """Descarga y cachea el logo Arauco (blanco PNG). Retorna None si falla."""
    global _ARAUCO_LOGO_BYTES
    if _ARAUCO_LOGO_BYTES is None:
        try:
            import requests as _req
            r = _req.get(_ARAUCO_LOGO_URL, timeout=5)
            if r.status_code == 200:
                _ARAUCO_LOGO_BYTES = r.content
            else:
                logging.warning("Logo Arauco: HTTP %s", r.status_code)
        except Exception as e:
            logging.warning("Logo Arauco no disponible: %s", e)
    return _ARAUCO_LOGO_BYTES


def trim_history(history: list) -> list:
    """Mantiene solo los últimos MAX_HISTORY mensajes (par usuario/asistente)."""
    return history[-MAX_HISTORY:]


def _safe_err(e: Exception) -> str:
    """Retorna mensaje de error seguro para mostrar al usuario sin exponer internals."""
    msg = str(e)
    # Redactar si hay credenciales en el mensaje
    for secret in ("api_key", "apikey", "token", "password", "bearer", "secret", "authorization"):
        if secret in msg.lower():
            return "Error interno del servidor."
    # Redactar URLs y rutas de archivo
    msg = re.sub(r"https?://\S+", "[URL]", msg)
    msg = re.sub(r"/[\w/.\-]+", "[ruta]", msg)
    return msg[:180]


def _cached_system(system: str) -> list:
    """Wraps a system prompt string with cache_control for prompt caching."""
    return [{"type": "text", "text": system, "cache_control": {"type": "ephemeral"}}]

# Pre-computado para evitar recrear el dict en cada llamada
_CACHED_SYSTEM_PROMPT: list | None = None


_EFFORT_MODELS = {"claude-sonnet-4-6", "claude-opus-4-6", "claude-opus-4-7"}


def claude_response(system: str, user_msg: str, max_tokens: int = 512,
                    model: str = "claude-haiku-4-5",
                    history: list | None = None,
                    effort: str = "low") -> str:
    """Llama a la API de Anthropic con historial conversacional opcional."""
    messages = list(history) if history else []
    messages.append({"role": "user", "content": user_msg})

    kwargs = dict(
        model=model,
        max_tokens=max_tokens,
        system=_cached_system(system),
        messages=messages,
        extra_headers={"anthropic-beta": "prompt-caching-2024-07-31"},
    )
    if model in _EFFORT_MODELS:
        kwargs["output_config"] = {"effort": effort}

    for attempt in range(3):
        try:
            response = client.messages.create(**kwargs)
            return response.content[0].text
        except anthropic.APIStatusError as e:
            if e.status_code == 500 and attempt < 2:
                time.sleep(2 ** attempt)
                continue
            raise


def push_history(context, user_msg: str, assistant_reply: str):
    """Agrega un turno al historial y lo recorta si es necesario."""
    history = context.user_data.setdefault("history", [])
    history.append({"role": "user",      "content": user_msg})
    history.append({"role": "assistant", "content": assistant_reply})
    context.user_data["history"] = trim_history(history)


ARTIFACT_KEYBOARD = InlineKeyboardMarkup([[
    InlineKeyboardButton("📁 Editables", callback_data="art_menu_editables"),
    InlineKeyboardButton("📋 Planner",   callback_data="art_menu_planner"),
    InlineKeyboardButton("📝 Nota",      callback_data="notas_modo"),
]])

EDITABLES_KEYBOARD = InlineKeyboardMarkup([[
    InlineKeyboardButton("📊 Excel", callback_data="art_excel"),
    InlineKeyboardButton("🖥️ PPT",  callback_data="art_pptx"),
    InlineKeyboardButton("📄 PDF",  callback_data="art_pdf"),
], [
    InlineKeyboardButton("📅 Gantt", callback_data="art_gantt"),
    InlineKeyboardButton("🌐 HTML",  callback_data="art_html"),
    InlineKeyboardButton("📧 Email", callback_data="art_email"),
], [
    InlineKeyboardButton("← Volver", callback_data="art_menu_back"),
]])

PLANNER_KEYBOARD = InlineKeyboardMarkup([[
    InlineKeyboardButton("🖥️ Desktop", callback_data="art_planner"),
    InlineKeyboardButton("📱 Mobile",  callback_data="art_planner_mobile"),
], [
    InlineKeyboardButton("← Volver",   callback_data="art_menu_back"),
]])

NOTAS_KEYBOARD = InlineKeyboardMarkup([[
    InlineKeyboardButton("📓 Juntar en OneNote", callback_data="notas_join"),
    InlineKeyboardButton("🗑 Borrar notas",      callback_data="notas_clear"),
], [
    InlineKeyboardButton("📍 Ubicación", callback_data="notas_add_location"),
], [
    InlineKeyboardButton("✅ Salir del modo notas", callback_data="notas_salir"),
]])

NOTAS_POST_KEYBOARD = InlineKeyboardMarkup([[
    InlineKeyboardButton("✏️ Editar",            callback_data="notas_editar"),
    InlineKeyboardButton("📄 → PDF",             callback_data="notas_pdf"),
], [
    InlineKeyboardButton("📓 → Word (OneNote)",  callback_data="notas_docx"),
    InlineKeyboardButton("🗑 Borrar notas",       callback_data="notas_clear"),
], [
    InlineKeyboardButton("✅ Salir del modo notas", callback_data="notas_salir"),
]])

IMAGE_PENDING_KEYBOARD = InlineKeyboardMarkup([[
    InlineKeyboardButton("🔍 Analizar imagen",   callback_data="img_analizar"),
], [
    InlineKeyboardButton("📊 Excel",  callback_data="img_art_excel"),
    InlineKeyboardButton("🌐 HTML",   callback_data="img_art_html"),
    InlineKeyboardButton("📄 PDF",    callback_data="img_art_pdf"),
], [
    InlineKeyboardButton("🖥️ PPT",    callback_data="img_art_pptx"),
    InlineKeyboardButton("📅 Gantt",  callback_data="img_art_gantt"),
    InlineKeyboardButton("📧 Email",  callback_data="img_art_email"),
]])

DOC_KEYBOARD = InlineKeyboardMarkup([[
    InlineKeyboardButton("📊 Excel", callback_data="art_excel"),
    InlineKeyboardButton("🖥️ PPT",   callback_data="art_pptx"),
    InlineKeyboardButton("📄 PDF",   callback_data="art_pdf"),
], [
    InlineKeyboardButton("📅 Gantt", callback_data="art_gantt"),
    InlineKeyboardButton("🌐 HTML",  callback_data="art_html"),
    InlineKeyboardButton("📧 Email", callback_data="art_email"),
], [
    InlineKeyboardButton("📥 Indexar RAG", callback_data="rag_index"),
]])

MODELS = {
    "haiku":  ("claude-haiku-4-5",  "⚡ Haiku",  "Rápido y económico"),
    "sonnet": ("claude-sonnet-4-6", "🧠 Sonnet", "Balanceado"),
    "opus":   ("claude-opus-4-7",   "🚀 Opus",   "Máxima capacidad"),
}
DEFAULT_MODEL = "haiku"

def get_model(context) -> str:
    """Retorna el model ID seleccionado por el usuario (o el default)."""
    key = context.user_data.get("model", DEFAULT_MODEL)
    return MODELS[key][0]

def get_model_label(context) -> str:
    key = context.user_data.get("model", DEFAULT_MODEL)
    return MODELS[key][1]

def get_model_for_msg(context, user_msg: str) -> str:
    """Usa Opus si el mensaje lo menciona explícitamente; si no, usa la preferencia del usuario."""
    if "opus" in user_msg.lower():
        return MODELS["opus"][0]
    return get_model(context)

def model_keyboard(current: str) -> InlineKeyboardMarkup:
    buttons = []
    for key, (_, label, desc) in MODELS.items():
        check = "✅ " if key == current else ""
        buttons.append([InlineKeyboardButton(f"{check}{label} — {desc}", callback_data=f"mdl_{key}")])
    return InlineKeyboardMarkup(buttons)

async def modelo_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    current = context.user_data.get("model", DEFAULT_MODEL)
    await update.message.reply_text(
        f"🤖 *Selecciona el modelo LLM*\nActual: {get_model_label(context)}",
        reply_markup=model_keyboard(current),
        parse_mode="Markdown"
    )

async def modelo_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    key = query.data.replace("mdl_", "")
    if key not in MODELS:
        return
    context.user_data["model"] = key
    _, label, desc = MODELS[key]
    try:
        await query.edit_message_text(
            f"✅ Modelo actualizado: *{label}*\n_{desc}_\n\nTodos los mensajes usarán este modelo.",
            reply_markup=model_keyboard(key),
            parse_mode="Markdown"
        )
    except Exception:
        pass


CODE_BLOCK_RE = re.compile(r'```(\w*)\n(.*?)```', re.DOTALL)
EXT_MAP = {
    "python": "py", "py": "py",
    "sql": "sql",
    "bash": "sh", "sh": "sh",
    "r": "r",
    "javascript": "js", "js": "js",
    "html": "html",
    "css": "css",
    "json": "json",
}
# Captions y nombres de archivo para tipos de bloque específicos
FILE_META = {
    "css":  ("estilos.css",  "🎨 Hoja de estilos"),
    "json": ("datos.json",   "📋 Datos en JSON"),
}

async def send_reply(update: Update, text: str, reply_markup=None, context=None):
    """
    Envía la respuesta de Claude a Telegram.
    Bloques de código ≥10 líneas se extraen y envían como archivo adjunto.
    Bloques HTML se envían como .html con caption descriptivo.
    """
    code_blocks = CODE_BLOCK_RE.findall(text)
    large_blocks = [(lang, code) for lang, code in code_blocks if code.count('\n') >= 10]

    if not large_blocks:
        # Sin código largo — envío normal con teclado de artefactos
        try:
            msg = await update.message.reply_text(
                fmt(text) + "\n\n🎨 <i>¿Generar un artefacto visual con esto?</i>",
                reply_markup=reply_markup,
                parse_mode="HTML"
            )
        except Exception:
            msg = await update.message.reply_text(
                text + "\n\n🎨 ¿Generar un artefacto visual con esto?",
                reply_markup=reply_markup
            )
        _track(context, msg.message_id)
        return

    # Hay bloques de código largos — enviar texto limpio primero
    clean_text = CODE_BLOCK_RE.sub("", text).strip()
    if clean_text:
        try:
            await update.message.reply_text(
                fmt(clean_text),
                reply_markup=reply_markup,
                parse_mode="HTML"
            )
        except Exception:
            await update.message.reply_text(clean_text, reply_markup=reply_markup)

    # Enviar cada bloque — HTML como URL, resto como archivo
    for i, (lang, code) in enumerate(large_blocks, 1):
        lang_key = lang.lower()
        if lang_key == "html":
            # HTML → servidor HTTP → URL que abre en browser externo
            url = store_html(code.strip())
            label = "Dashboard" if i == 1 else f"Artefacto {i}"
            await update.message.reply_text(
                f"🌐 <b>{label} HTML listo</b>\n\nToca el enlace para abrirlo en tu browser:\n{url}",
                parse_mode="HTML"
            )
        else:
            ext = EXT_MAP.get(lang_key, "txt")
            if lang_key in FILE_META:
                filename, caption = FILE_META[lang_key]
                if len(large_blocks) > 1:
                    base, dot_ext = filename.rsplit(".", 1)
                    filename = f"{base}_{i}.{dot_ext}"
            else:
                filename = f"script_{i}.{ext}" if len(large_blocks) > 1 else f"script.{ext}"
                caption  = f"📎 {filename} — copia y ejecuta en tu entorno"
            buf = io.BytesIO(code.strip().encode("utf-8"))
            await update.message.reply_document(document=buf, filename=filename, caption=caption)


# Mapa de palabras clave → tipo de artefacto
_ARTIFACT_INTENT = {
    "html":         ["dashboard", "html", "interactivo", "visualización", "visualizacion"],
    "excel":        ["excel", "tabla", "spreadsheet", "hoja de cálculo", "hoja de calculo"],
    "pdf":          ["pdf", "informe", "reporte", "report"],
    "gantt":        ["gantt", "cronograma", "carta gantt", "timeline", "plan de proyecto"],
    "pptx":         ["ppt", "pptx", "powerpoint", "presentación", "presentacion", "diapositiva"],
    "email":        ["envía un correo", "envia un correo", "manda un correo", "redacta un correo", "escribe un correo"],
    "planner":        ["planner mc", "dashboard planner", "tablero mc", "dashboard mc"],
    "planner_mobile": ["planner mobile", "planner celular", "dashboard mobile", "tablero mobile"],
}

def _detect_artifact_intent(text: str) -> str | None:
    lower = text.lower()
    for artifact_type, keywords in _ARTIFACT_INTENT.items():
        for kw in keywords:
            if re.search(r'\b' + re.escape(kw) + r'\b', lower):
                return artifact_type
    return None


def _user_msg_has_description(user_msg: str) -> bool:
    """True si el mensaje contiene una descripción propia más allá del keyword del artefacto."""
    lower = user_msg.lower()
    # Quitar keywords de artefactos y palabras de comando genéricas
    stopwords = ["haz", "hacer", "genera", "generar", "crea", "crear", "dame", "quiero",
                 "un", "una", "el", "la", "por", "favor", "porfavor", "please",
                 "html", "excel", "pdf", "gantt", "pptx", "ppt", "email", "dashboard",
                 "informe", "reporte", "presentación", "presentacion", "tabla", "correo"]
    tokens = [t for t in lower.split() if t not in stopwords]
    return len(tokens) >= 3  # hay descripción real si quedan ≥3 palabras significativas


def _build_artifact_description(user_msg: str, context) -> str:
    """Construye la descripción para el artefacto.

    Si el usuario escribió una descripción propia (más allá del keyword),
    la usa sin inyectar el documento — el usuario está pidiendo algo diferente.
    Si solo presionó un botón o usó un keyword simple, enriquece con el contexto del doc.
    """
    last_analysis   = context.user_data.get("last_analysis", "")
    structured_data = context.user_data.get("structured_data", {})
    doc_content     = context.user_data.get("doc_content", "")
    doc_tipo        = context.user_data.get("doc_tipo", "")

    base = user_msg or "Basado en el análisis previo"

    # Si el usuario describió algo específico, usar solo su descripción — sin heredar contexto anterior
    if user_msg and _user_msg_has_description(user_msg):
        return base

    # Sin descripción propia (botón o keyword simple): usa el doc/análisis completo
    if structured_data:
        data_block = json.dumps(structured_data, ensure_ascii=False, default=str)
        return f"{base}\n\nDATOS EXACTOS DEL ARCHIVO {doc_tipo} (úsalos LITERALMENTE):\n{data_block[:8000]}"
    elif doc_content:
        return f"{base}\n\nCONTENIDO COMPLETO DEL ARCHIVO {doc_tipo}:\n{doc_content[:8000]}"
    elif last_analysis:
        return f"{base}\n\nContexto del análisis previo:\n{last_analysis}"
    return base


async def _render_artifact(artifact_type: str, description: str,
                           reply_fn, context) -> None:
    """Genera y envía un artefacto, luego muestra el teclado para generar otro."""
    # Planner MC desktop: HTML estático, sin llamada a Claude
    if artifact_type == "planner":
        if not _PLANNER_HTML:
            await reply_fn("⚠️ Dashboard Planner no disponible en este entorno.")
            return
        url = store_html(_PLANNER_HTML)
        await reply_fn(
            f"📋 <b>Dashboard Planner MC listo</b>\n\nToca el enlace para abrirlo:\n{url}",
            parse_mode="HTML", reply_markup=ARTIFACT_KEYBOARD
        )
        return

    # Planner MC mobile: versión optimizada para celular
    if artifact_type == "planner_mobile":
        if not _PLANNER_MOBILE_HTML:
            await reply_fn("⚠️ Dashboard Mobile no disponible en este entorno.")
            return
        url = store_html(_PLANNER_MOBILE_HTML)
        await reply_fn(
            f"📱 <b>Planner MC Mobile listo</b>\n\nToca el enlace para abrirlo en tu celular:\n{url}",
            parse_mode="HTML", reply_markup=ARTIFACT_KEYBOARD
        )
        return

    _tokens_map = {"html": 8000, "pdf": 6000, "gantt": 4000, "excel": 3000, "pptx": 6000, "email": 2000, "notas_onenote": 8000}
    prompt = ARTIFACT_PROMPTS[artifact_type].replace("{CSS_URL}", f"{PUBLIC_BASE}/arauco.css")
    raw = claude_response(prompt, description,
                          max_tokens=_tokens_map.get(artifact_type, 4000),
                          model="claude-sonnet-4-6",
                          effort="high")
    try:
        if artifact_type in ("html", "notas_onenote"):
            html = raw.strip()
            if html.startswith("```"):
                html = html.split("\n", 1)[-1]
            if html.endswith("```"):
                html = html.rsplit("```", 1)[0]
            html = html.strip()
            if not html.lower().startswith("<!doctype") and "<html" not in html.lower():
                await reply_fn("⚠️ El HTML generado está incompleto. Intenta de nuevo.")
                return
            url = store_html(html)
            if artifact_type == "notas_onenote":
                context.user_data["last_notas_data"] = context.user_data.get("notas", [])
                await reply_fn(f"📓 <b>Documento listo</b>\n\nToca el enlace para abrirlo:\n{url}",
                               parse_mode="HTML", reply_markup=NOTAS_POST_KEYBOARD)
            else:
                await reply_fn(f"🌲 <b>Dashboard interactivo listo</b>\n\nToca el enlace:\n{url}",
                               parse_mode="HTML", reply_markup=ARTIFACT_KEYBOARD)
            return
        elif artifact_type == "excel":
            data = extract_json(raw)
            buf  = build_excel(data)
            filename = f"arauco-{data.get('titulo','datos').lower().replace(' ','-')[:30]}.xlsx"
            await reply_fn(buf=buf, filename=filename, caption="📊 Excel generado — Arauco Mejora Continua")
        elif artifact_type == "pdf":
            data  = extract_json(raw)
            buf   = build_pdf(data)
            titulo = data.get("titulo", "informe-arauco").lower().replace(" ", "-")[:30]
            await reply_fn(buf=buf, filename=f"{titulo}.pdf", caption="📄 Informe PDF — Arauco Mejora Continua")
        elif artifact_type == "gantt":
            data  = extract_json(raw)
            url   = build_gantt(data)
            titulo = data.get("titulo", "Carta Gantt")
            await reply_fn(f"📅 <b>{titulo}</b>\n\nToca el enlace:\n{url}",
                           parse_mode="HTML", reply_markup=ARTIFACT_KEYBOARD)
            return
        elif artifact_type == "pptx":
            data  = extract_json(raw)
            buf   = build_pptx(data)
            titulo = data.get("titulo", "presentacion-arauco").lower().replace(" ", "-")[:30]
            await reply_fn(buf=buf, filename=f"{titulo}.pptx", caption="🖥️ PowerPoint generado — Arauco Mejora Continua")
        elif artifact_type == "email":
            data = extract_json(raw)
            context.user_data["pending_email"] = data
            context.user_data["waiting_email_recipient"] = True
            await reply_fn("📧 Borrador listo. ¿A qué correo lo envío?")
            return
        # Para archivos (excel, pdf, pptx): muestra teclado en mensaje separado
        await reply_fn("¿Generar otro artefacto?", reply_markup=ARTIFACT_KEYBOARD)
    except (json.JSONDecodeError, ValueError) as e:
        await reply_fn(f"⚠️ Error al parsear respuesta: {str(e)[:120]}\nInicio: {raw[:200]}")
    except Exception as e:
        await reply_fn(f"⚠️ Error generando artefacto: {_safe_err(e)}")


def _make_reply_fn(message):
    async def reply_fn(text=None, *, buf=None, filename=None, caption=None, **kwargs):
        if buf is not None:
            await message.reply_document(document=buf, filename=filename, caption=caption)
        else:
            await message.reply_text(text, **kwargs)
    return reply_fn


def _email_preview(data: dict, n_adjuntos: int = 0) -> str:
    lines = [
        f"📧 <b>Borrador de correo</b>\n",
        f"<b>Para:</b> {_html.escape(data.get('para', ''))}",
    ]
    if data.get("cc"):
        lines.append(f"<b>CC:</b> {_html.escape(data['cc'])}")
    lines.append(f"<b>Asunto:</b> {_html.escape(data.get('asunto', ''))}\n")
    lines.append(_html.escape(data.get("cuerpo", "")))
    if n_adjuntos:
        lines.append(f"\n📎 <i>{n_adjuntos} imagen(es) adjunta(s)</i>")
    return "\n".join(lines)


async def _handle_nlm_query(update: Update, context: ContextTypes.DEFAULT_TYPE, question: str):
    """Responde pregunta via RAG + Claude con formato estructurado."""
    chunks = rag.query(question)
    if not chunks:
        await update.message.reply_text(
            f"🔍 Sin resultados para <i>{_html.escape(question)}</i> en la base de conocimiento.\n\n"
            "Los documentos están indexados pero ningún fragmento superó el umbral de relevancia. "
            "Intenta reformular la pregunta o usa /notebookrag para ver los documentos disponibles.",
            parse_mode="HTML"
        )
        return

    rag_context = rag.build_context(question)
    system = (
        "Eres el asistente de base de conocimiento de Arauco — Subgerencia de Mejora Continua.\n"
        "Responde ÚNICAMENTE basado en los documentos del contexto.\n"
        "Si la información no está en los documentos, indícalo claramente.\n\n"
        "FORMATO — adapta según la complejidad de la pregunta:\n\n"
        "- Encabezado: ━━ 📋 [título breve] ━━\n"
        "- Párrafo introductorio si la respuesta lo amerita\n"
        "- Secciones numeradas con emoji (1️⃣ 2️⃣ ...) solo si hay varios temas distintos\n"
        "- Bullets 🔹 para listas de puntos\n"
        "- **Negrita** para términos técnicos clave y valores críticos\n"
        "- (FUENTE: archivo) al citar datos específicos del documento\n"
        "- ⚠️ **Nota:** al final solo si hay limitaciones o información faltante\n"
        "- Para preguntas simples, respuesta directa sin secciones"
        + rag_context
    )

    try:
        resp = claude_response(system, question, max_tokens=1200)
        await update.message.reply_text(fmt(resp), parse_mode="HTML")
    except Exception as e:
        await update.message.reply_text(f"⚠️ Error: {_safe_err(e)}")


async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_msg = update.message.text

    # Intercepta modo NotebookRAG
    if context.user_data.get("nlm_mode"):
        context.user_data.pop("nlm_mode")
        await _handle_nlm_query(update, context, user_msg)
        return

    # Intercepta edición de campo del email
    if context.user_data.get("editing_email_field"):
        campo = context.user_data.pop("editing_email_field")
        draft = context.user_data.get("pending_email", {})
        draft[campo] = user_msg.strip()
        context.user_data["pending_email"] = draft
        n_adj = context.user_data.get("email_n_adjuntos", 0)
        await update.message.reply_text(
            _email_preview(draft, n_adj), parse_mode="HTML",
            reply_markup=EMAIL_CONFIRM_KEYBOARD
        )
        return

    # Intercepta respuesta con destinatario de email
    if context.user_data.get("waiting_email_recipient"):
        context.user_data.pop("waiting_email_recipient")
        recipient = user_msg.strip()
        if not re.fullmatch(r"[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}", recipient):
            await update.message.reply_text("⚠️ El correo ingresado no parece válido. Intenta de nuevo.")
            context.user_data["waiting_email_recipient"] = True
            return
        draft = context.user_data.get("pending_email", {})
        draft["para"] = recipient
        context.user_data["pending_email"] = draft
        n_adj = context.user_data.get("email_n_adjuntos", 0)
        await update.message.reply_text(
            _email_preview(draft, n_adj), parse_mode="HTML",
            reply_markup=EMAIL_CONFIRM_KEYBOARD
        )
        return

    # ── Modo edición de notas: instrucciones para re-generar el HTML ──
    if context.user_data.get("modo_editar_notas"):
        context.user_data.pop("modo_editar_notas", None)
        notas_txt = context.user_data.get("last_notas_txt", "")
        if not notas_txt:
            await update.message.reply_text("⚠️ No hay documento base para editar.")
            return
        description = (
            f"[INSTRUCCIÓN DE EDICIÓN]: {user_msg}\n\n"
            f"[NOTAS ORIGINALES]:\n{notas_txt}"
        )
        await update.message.reply_text("⏳ Aplicando edición al documento...")
        await _render_artifact("notas_onenote", description, _make_reply_fn(update.message), context)
        return

    # ── Modo notas: cualquier mensaje se guarda como nota ──
    _nota_kw = ("nota:", "nota ", "apunta:", "apunta ", "anota:", "anota ", "registra:", "registra ")
    _es_nota_explicita = user_msg.lower().startswith(_nota_kw)
    if context.user_data.get("modo_notas") or _es_nota_explicita:
        if _es_nota_explicita:
            sep = ":" if ":" in user_msg.split()[0] else " "
            content = user_msg.split(sep, 1)[-1].strip()
        else:
            content = user_msg.strip()
        if content:
            notas = context.user_data.get("notas", [])
            notas.append({"texto": content, "fecha": datetime.now().strftime("%d/%m %H:%M"), "n": len(notas) + 1})
            context.user_data["notas"] = notas
            await update.message.reply_text(
                _notas_status_text(notas), parse_mode="Markdown", reply_markup=NOTAS_KEYBOARD
            )
            return

    # Detecta intención de artefacto y genera directamente usando todo el contexto disponible
    artifact_intent = _detect_artifact_intent(user_msg)
    if artifact_intent and artifact_intent in ARTIFACT_PROMPTS:
        description = _build_artifact_description(user_msg, context)
        await update.message.reply_text(f"⏳ Generando *{artifact_intent}*...", parse_mode="Markdown")
        await _render_artifact(artifact_intent, description,
                               _make_reply_fn(update.message), context)
        return

    try:
        history    = context.user_data.get("history", [])
        try:
            rag_ctx = rag.build_context(user_msg) if rag.col and rag.col.count() > 0 else ""
        except Exception:
            rag_ctx = ""
        system     = SYSTEM_PROMPT + rag_ctx
        reply = claude_response(system, user_msg, model=get_model_for_msg(context, user_msg), history=history)
        push_history(context, user_msg, reply)
        context.user_data["last_analysis"] = reply
        await send_reply(update, reply, reply_markup=ARTIFACT_KEYBOARD, context=context)
    except Exception as e:
        await update.message.reply_text(f"⚠️ Error al procesar: {_safe_err(e)}")


async def handle_image(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not update.message or not update.message.photo:
        return
    photo = update.message.photo[-1]
    file = await context.bot.get_file(photo.file_id)
    try:
        file_bytes = await asyncio.wait_for(file.download_as_bytearray(), timeout=30.0)
    except asyncio.TimeoutError:
        await update.message.reply_text("⚠️ La descarga tardó demasiado. Intenta con una imagen más pequeña.")
        return
    image_b64 = base64.standard_b64encode(file_bytes).decode("utf-8")
    caption = update.message.caption or ""

    images = context.user_data.get("pending_images", [])
    images.append({"b64": image_b64, "media_type": "image/jpeg", "caption": caption})
    context.user_data["pending_images"] = images
    context.user_data["pending_image"] = images[0]

    user_id = update.effective_user.id
    message  = update.message

    # Cancela el timer anterior y reinicia (debounce 1.5 s)
    prev = _img_debounce.pop(user_id, None)
    if prev and not prev.done():
        prev.cancel()

    async def _show_keyboard():
        try:
            await asyncio.sleep(1.5)
            n = len(context.user_data.get("pending_images", []))
            label = "imagen recibida" if n == 1 else f"{n} imágenes recibidas"
            await message.reply_text(
                f"🖼 *{label.capitalize()}.*\n\n¿Qué hacemos con {'ella' if n == 1 else 'ellas'}?",
                parse_mode="Markdown",
                reply_markup=IMAGE_PENDING_KEYBOARD,
            )
        except asyncio.CancelledError:
            return
        finally:
            # Solo limpia si este task sigue siendo el activo; evita que tareas
            # canceladas anteriores eliminen el task más reciente del dict.
            if _img_debounce.get(user_id) is asyncio.current_task():
                _img_debounce.pop(user_id, None)

    _img_debounce[user_id] = asyncio.create_task(_show_keyboard())


async def _analyze_image(context) -> str | None:
    """Llama a Claude con todas las imágenes pendientes; guarda y devuelve el análisis."""
    images = context.user_data.get("pending_images") or []
    if not images:
        img = context.user_data.get("pending_image")
        if not img:
            return None
        images = [img]

    caption = next((i["caption"] for i in images if i.get("caption")), "") or (
        "Analiza estas imágenes en el contexto operacional forestal de Arauco. "
        "Identifica equipos, procesos, problemas o métricas relevantes."
    )
    content = [
        {"type": "image", "source": {"type": "base64",
                                     "media_type": i["media_type"],
                                     "data": i["b64"]}}
        for i in images
    ]
    content.append({"type": "text", "text": caption})

    response = client.messages.create(
        model=get_model(context),
        max_tokens=1024,
        system=_cached_system(SYSTEM_PROMPT),
        messages=[{"role": "user", "content": content}],
        extra_headers={"anthropic-beta": "prompt-caching-2024-07-31"},
    )
    if not response.content:
        return None
    analysis = response.content[0].text
    push_history(context, caption, analysis)
    context.user_data["last_analysis"] = analysis
    return analysis


async def image_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()

    if not context.user_data.get("pending_image"):
        await query.edit_message_text("⚠️ No hay imagen pendiente.")
        return

    if query.data == "img_analizar":
        await query.edit_message_text("⏳ Analizando imagen...")
        analysis = await _analyze_image(context)
        context.user_data.pop("pending_images", None)
        if not analysis:
            await query.message.reply_text("⚠️ No se pudo analizar la imagen.")
            return
        try:
            await query.message.reply_text(
                fmt(analysis) + "\n\n🎨 <i>¿Generar un artefacto visual con esto?</i>",
                parse_mode="HTML", reply_markup=ARTIFACT_KEYBOARD,
            )
        except Exception:
            await query.message.reply_text(
                analysis + "\n\n🎨 ¿Generar un artefacto visual con esto?",
                reply_markup=ARTIFACT_KEYBOARD,
            )
        return

    # img_art_<tipo>
    art_type = query.data.replace("img_art_", "")

    # ── PDF con imágenes reales ───────────────────────────────────────
    if art_type == "pdf":
        images = list(context.user_data.get("pending_images", []))
        context.user_data.pop("pending_images", None)
        n = len(images)
        await query.edit_message_text(
            f"⏳ Analizando {n} imagen(es) y generando PDF...",
            parse_mode="Markdown"
        )
        pdf_prompt = (
            "Analiza estas imágenes en el contexto operacional forestal de Arauco "
            "y devuelve ÚNICAMENTE un JSON válido con esta estructura:\n"
            '{"titulo":"...","subtitulo":"...","area":"Subgerencia de Mejora Continua",'
            f'"fecha":"{datetime.now().strftime("%d de %B de %Y")}",'
            '"imagenes":['
            '{"titulo":"...","descripcion":"...","hallazgos":["...","..."]}'
            f'],  // exactamente {n} objeto(s)\n'
            '"conclusiones":["...","..."]}'
        )
        content = [
            {"type": "image", "source": {"type": "base64",
                                         "media_type": i["media_type"],
                                         "data": i["b64"]}}
            for i in images
        ]
        content.append({"type": "text", "text": pdf_prompt})
        resp = client.messages.create(
            model=get_model(context), max_tokens=3000,
            system=_cached_system(SYSTEM_PROMPT),
            messages=[{"role": "user", "content": content}],
            extra_headers={"anthropic-beta": "prompt-caching-2024-07-31"},
        )
        try:
            if not resp.content:
                raise ValueError("Respuesta vacía del modelo")
            pdf_data = extract_json(resp.content[0].text)
            buf = build_pdf_imagenes(pdf_data, images)
            titulo = pdf_data.get("titulo", "informe-visual")[:30].lower().replace(" ", "-")
            await query.message.reply_document(
                document=buf,
                filename=f"{titulo}.pdf",
                caption="📄 Informe PDF con imágenes — Arauco Mejora Continua",
                reply_markup=ARTIFACT_KEYBOARD,
            )
        except Exception as e:
            await query.message.reply_text(f"⚠️ Error generando PDF: {_safe_err(e)}")
        return

    # ── PPT con imágenes reales ───────────────────────────────────────
    if art_type == "pptx":
        images = list(context.user_data.get("pending_images", []))
        context.user_data.pop("pending_images", None)
        n = len(images)
        await query.edit_message_text(
            f"⏳ Analizando {n} imagen(es) y generando presentación...",
            parse_mode="Markdown"
        )
        pptx_prompt = (
            "Analiza estas imágenes en el contexto operacional forestal de Arauco "
            "y devuelve ÚNICAMENTE un JSON válido con esta estructura:\n"
            '{"titulo":"...","subtitulo":"...","area":"Subgerencia de Mejora Continua",'
            f'"fecha":"{datetime.now().strftime("%d de %B de %Y")}",'
            '"autor":"Arauco","imagenes":['
            '{"titulo":"...","descripcion":"...","hallazgos":["...","..."]}'
            f'],  // exactamente {n} objeto(s)\n'
            '"conclusiones":["...","..."]}'
        )
        content = [
            {"type": "image", "source": {"type": "base64",
                                         "media_type": i["media_type"],
                                         "data": i["b64"]}}
            for i in images
        ]
        content.append({"type": "text", "text": pptx_prompt})
        resp = client.messages.create(
            model=get_model(context), max_tokens=3000,
            system=_cached_system(SYSTEM_PROMPT),
            messages=[{"role": "user", "content": content}],
            extra_headers={"anthropic-beta": "prompt-caching-2024-07-31"},
        )
        try:
            if not resp.content:
                raise ValueError("Respuesta vacía del modelo")
            pptx_data = extract_json(resp.content[0].text)
            buf = build_pptx_imagenes(pptx_data, images)
            titulo = pptx_data.get("titulo", "presentacion-arauco")[:30].lower().replace(" ", "-")
            await query.message.reply_document(
                document=buf,
                filename=f"{titulo}.pptx",
                caption="🖥️ PowerPoint con imágenes — Arauco Mejora Continua",
                reply_markup=ARTIFACT_KEYBOARD,
            )
        except Exception as e:
            await query.message.reply_text(f"⚠️ Error generando PPT: {_safe_err(e)}")
        return

    if art_type not in ARTIFACT_PROMPTS:
        await query.answer("Tipo de artefacto no reconocido.", show_alert=True)
        return

    await query.edit_message_text(f"⏳ Analizando imagen y generando *{art_type}*...",
                                  parse_mode="Markdown")
    analysis = await _analyze_image(context)
    if art_type == "email":
        images = context.user_data.get("pending_images", [])
        context.user_data["email_attachments"] = list(images)
        context.user_data["email_n_adjuntos"] = len(images)
    context.user_data.pop("pending_images", None)
    if not analysis:
        await query.message.reply_text("⚠️ No se pudo analizar la imagen.")
        return
    await _render_artifact(art_type, analysis, _make_reply_fn(query.message), context)


def _notas_status_text(notas: list) -> str:
    if not notas:
        return "📝 *Modo notas* — escribe o envía un audio."
    n = len(notas)
    preview = notas[-1]["texto"][:60] + ("…" if len(notas[-1]["texto"]) > 60 else "")
    return (f"📝 *Nota {n} guardada*\n"
            f"_{preview}_\n\n"
            f"Modo notas activo — escribe o envía un audio.")


async def notas_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()

    if query.data == "notas_modo":
        context.user_data["modo_notas"] = True
        notas = context.user_data.get("notas", [])
        n = len(notas)
        msg = ("📝 *Modo notas activado*\n\n"
               "Escribe lo que quieras anotar o envía un audio — "
               "todo se guardará como nota.\n\n"
               f"{'_Tienes ' + str(n) + ' nota(s) previas._' if n else '_Aún no hay notas._'}")
        try:
            await query.edit_message_reply_markup(reply_markup=None)
        except Exception:
            pass
        await query.message.reply_text(msg, parse_mode="Markdown", reply_markup=NOTAS_KEYBOARD)
        return

    if query.data == "notas_salir":
        context.user_data.pop("modo_notas", None)
        notas = context.user_data.get("notas", [])
        n = len(notas)
        await query.edit_message_text(
            f"✅ Modo notas desactivado.{' Tienes ' + str(n) + ' nota(s) guardada(s).' if n else ''}",
            reply_markup=ARTIFACT_KEYBOARD if n else None
        )
        return

    if query.data == "notas_add_location":
        notas = context.user_data.get("notas", [])
        if not notas:
            await query.answer("⚠️ Escribe una nota primero.", show_alert=True)
            return
        context.user_data["agregar_ubicacion_a_nota"] = len(notas)  # índice (1-based) de la nota target
        await query.message.reply_text(
            "📍 Comparte tu ubicación para añadirla a la última nota.",
            parse_mode="Markdown"
        )
        return

    if query.data == "notas_clear":
        context.user_data.pop("notas", None)
        context.user_data.pop("modo_notas", None)
        await query.edit_message_text("🗑 Notas borradas.")
        return

    if query.data == "notas_editar":
        notas = context.user_data.get("last_notas_data") or context.user_data.get("notas", [])
        if not notas:
            await query.answer("⚠️ No hay documento generado para editar.", show_alert=True)
            return
        context.user_data["modo_editar_notas"] = True
        await query.message.reply_text(
            "✏️ *Modo edición activo*\n\n"
            "Escribe las instrucciones de edición, por ejemplo:\n"
            "• _«agrega una sección de conclusiones»_\n"
            "• _«nota 3: nuevo texto corregido»_\n"
            "• _«cambia el título a Reunión Q2»_",
            parse_mode="Markdown"
        )
        return

    if query.data == "notas_pdf":
        notas = context.user_data.get("last_notas_data") or context.user_data.get("notas", [])
        if not notas:
            await query.answer("⚠️ No hay notas para exportar.", show_alert=True)
            return
        await query.message.reply_text("⏳ Generando PDF...")
        titulo = f"Notas — {datetime.now().strftime('%d/%m/%Y')}"
        buf = build_notas_pdf(notas, titulo)
        filename = f"notas-arauco-{datetime.now().strftime('%Y%m%d')}.pdf"
        await query.message.reply_document(document=buf, filename=filename,
                                           caption="📄 Notas en PDF — Arauco Mejora Continua")
        return

    if query.data == "notas_docx":
        notas = context.user_data.get("last_notas_data") or context.user_data.get("notas", [])
        if not notas:
            await query.answer("⚠️ No hay notas para exportar.", show_alert=True)
            return
        await query.message.reply_text("⏳ Generando Word...")
        titulo = f"Notas — {datetime.now().strftime('%d/%m/%Y')}"
        buf = build_notas_docx(notas, titulo)
        filename = f"notas-arauco-{datetime.now().strftime('%Y%m%d')}.docx"
        await query.message.reply_document(document=buf, filename=filename,
                                           caption="📓 Notas en Word — importa este archivo en OneNote")
        return

    if query.data != "notas_join":
        await query.answer("Acción no reconocida.", show_alert=True)
        return

    notas = context.user_data.get("notas", [])
    if not notas:
        await query.edit_message_text("⚠️ No hay notas guardadas.")
        return

    context.user_data.pop("modo_notas", None)
    await query.edit_message_reply_markup(reply_markup=None)
    await query.message.reply_text(f"⏳ Organizando *{len(notas)} notas* en documento OneNote...", parse_mode="Markdown")

    notas_txt = "\n\n".join([f"[Nota {n['n']} — {n['fecha']}]\n{n['texto']}" for n in notas])
    context.user_data["last_notas_txt"] = notas_txt
    description = f"{len(notas)} notas para organizar:\n\n{notas_txt}"

    await _render_artifact("notas_onenote", description, _make_reply_fn(query.message), context)


async def artifact_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()

    artifact_type = query.data.replace("art_", "")

    # Navegación de menús
    if artifact_type == "menu_editables":
        await query.edit_message_reply_markup(reply_markup=EDITABLES_KEYBOARD)
        return
    if artifact_type == "menu_planner":
        await query.edit_message_reply_markup(reply_markup=PLANNER_KEYBOARD)
        return
    if artifact_type == "menu_back":
        await query.edit_message_reply_markup(reply_markup=ARTIFACT_KEYBOARD)
        return

    last_analysis = context.user_data.get("last_analysis", "")

    _static_types = {"planner", "planner_mobile"}
    if not last_analysis and artifact_type not in _static_types:
        await query.edit_message_text("⚠️ No hay análisis previo. Envía una imagen o un mensaje primero.")
        return

    await query.edit_message_reply_markup(reply_markup=None)
    await query.message.reply_text(f"⏳ Generando *{artifact_type}*...", parse_mode="Markdown")

    description = _build_artifact_description("", context)
    await _render_artifact(artifact_type, description, _make_reply_fn(query.message), context)


async def skill_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    command = update.message.text.split()[0].lstrip("/")
    args = " ".join(context.args) if context.args else ""

    ejemplos = {
        "spec": "/spec alertas de falla cosechadoras Tigercat predio remoto",
        "plan": "/plan telemetría transporte rollizos temporada 2026",
        "build": "/build pipeline OEE equipos cosecha desde API John Deere",
        "test": "/test modelo predictivo procesadoras línea sur",
        "review": "/review propuesta optimización flota Opti-Cliente",
        "ship": "/ship dashboard diario avance cosecha para supervisores",
    }

    if not args:
        await update.message.reply_text(
            f"{SKILL_PROMPTS[command]}\n\n📌 *Ejemplo:* `{ejemplos[command]}`",
            parse_mode="Markdown"
        )
        return

    skill_system = SYSTEM_PROMPT + "\n\n" + SKILL_PROMPTS[command]
    reply = claude_response(skill_system, args, max_tokens=800, model=get_model(context))
    await update.message.reply_text(fmt(reply), parse_mode="HTML")


ARTIFACT_HELP = """🎨 */artifact* — Genera un archivo y lo envía aquí

*Tipos disponibles:*
• `html`  — Dashboard interactivo
• `excel` — Tabla de datos en formato Excel
• `pdf`   — Informe ejecutivo en PDF
• `gantt` — Carta Gantt del proyecto
• `pptx`  — Presentación PowerPoint
• `email` — Borrador de correo Outlook

*Uso:*
`/artifact html dashboard OEE semanal línea 3`
`/artifact excel tabla KPIs cosecha por turno`
`/artifact pdf informe pérdidas semana 23`
`/artifact gantt proyecto mejora bomba 42`
`/artifact pptx presentación resultados Q2 cosecha`
`/artifact email resumen análisis para juan@arauco.com`"""

ARTIFACT_PROMPTS = {
    "planner":        "",  # HTML estático desktop
    "planner_mobile": "",  # HTML estático mobile
    "notas_onenote": """Eres un asistente que organiza notas en un documento HTML estilo OneNote, responsive para desktop Y mobile.

El usuario te entrega N notas en texto libre. Tu tarea:
1. Leer todas las notas y detectar temas en común
2. Agrupar las notas por tema (si hay temas claros) o mantenerlas cronológicas
3. Generar un HTML completo, autocontenido y 100% responsive

ESTRUCTURA HTML obligatoria:
- <meta name="viewport" content="width=device-width, initial-scale=1.0, viewport-fit=cover">
- Fuente: DM Sans via Google Fonts (wght@400;500;600;700)
- Colores corporativos Arauco: --gris:#696158  --verde:#BFB800  --naranja:#EA7600  --fondo:#f5f4f0

DESKTOP (min-width: 768px):
- Layout flex: sidebar izquierdo 200px (fondo #696158, color blanco) + contenido principal
- Sidebar: logo `<img src="https://arauco.com/chile/wp-content/themes/arauco/assets/img/logo-arauco-blanco.png" alt="Arauco" height="24" style="display:block;margin:0 auto 12px">`, nombre del cuaderno, lista de secciones clicables con punto de color
- Contenido: fondo blanco, padding 40px, sombra sutil, max-width 760px

MOBILE (max-width: 767px):
- Sin sidebar — oculto con display:none
- Header fijo en top: fondo #696158, logo `<img src="https://arauco.com/chile/wp-content/themes/arauco/assets/img/logo-arauco-blanco.png" alt="Arauco" height="20" style="vertical-align:middle;margin-right:8px">` + título del documento, botón ☰ que abre drawer lateral
- Drawer: panel lateral deslizable (transform translateX) con las secciones
- Overlay oscuro al abrir drawer
- Padding seguro: padding-bottom: env(safe-area-inset-bottom, 16px)
- Botón flotante "📋 Copiar" fijo en bottom-right

CONTENIDO de cada sección:
- Título de sección: negrita, borde izquierdo 3px con color único por sección
- Cada nota: timestamp pequeño arriba en gris, texto en párrafo, separador sutil entre notas
- Sección final "⚠️ Pendientes" si hay tareas o acciones detectadas en el texto

BOTÓN COPIAR (ambas versiones):
- Copia texto plano con estructura: título, secciones, notas y timestamps
- Usa navigator.clipboard.writeText()
- Toast de confirmación al copiar

IMPORTANTE:
- Infiere un título relevante para el documento según el contenido
- Usa colores distintos por sección (verde oliva, naranja, gris tierra, azul, etc.)
- Responde SOLO con el HTML completo, sin texto adicional, sin markdown
""",
    "html": """Eres el Agente DA de Arauco — Subgerencia de Mejora Continua.
Genera un dashboard HTML interactivo y autocontenido.

HEAD obligatorio:
<link rel="stylesheet" href="{CSS_URL}">
<link href="https://fonts.googleapis.com/css2?family=Lato:wght@300;400;700;900&display=swap" rel="stylesheet">
<script src="https://cdn.jsdelivr.net/npm/chart.js"></script>

ESTRUCTURA:
1. Header — .dashboard-header con logo `<img src="https://arauco.com/chile/wp-content/themes/arauco/assets/img/logo-arauco-blanco.png" alt="Arauco" height="32" style="margin-right:16px">` + .dashboard-title + .dashboard-subtitle (fuente + fecha)
2. KPI cards — .grid.grid-4, mínimo 4, con .kpi-label / .kpi-value / .kpi-change (.positive/.negative/.neutral)
3. Filtros — .filtros-bar con <select id="f-COL"> por columna categórica + .btn-limpiar + .conteo-badge
4. <script> completo con: const DATOS=[...], const FILTROS_COLS=[...], const charts={}, funciones aplicarFiltros() / renderTabla() / actualizarGraficos() / limpiarFiltros(), inicialización en DOMContentLoaded
5. Mínimo 2 gráficos Chart.js (canvas) en .grid.grid-2 dentro de .card
6. Tabla filtrable en .card con overflow-x:auto (máx 50 filas)
7. Footer — .dashboard-footer

COLORES JS: grisTierra:'#696158', verdeOliva:'#BFB800', naranja:'#EA7600'

DATOS:
- EXCEL → usa stats[col].frecuencias para gráficos; muestra_top20 para DATOS[]
- PDF/WORD → extrae tablas y cifras del texto
- Nunca inventes cifras; formato chileno 1.234,5; .badge-null para vacíos

Responde ÚNICAMENTE con HTML. Sin markdown. Empieza con <!DOCTYPE html>.""",

    "excel": """Genera datos estructurados en formato JSON para crear un archivo Excel.

El JSON debe tener exactamente esta estructura:
{
  "titulo": "Nombre de la hoja",
  "encabezados": ["Col1", "Col2", "Col3", ...],
  "filas": [
    ["valor1", "valor2", "valor3", ...],
    ...
  ]
}

Los datos deben ser realistas para el contexto forestal de Arauco pedido (KPIs, mérdidas, productividad, equipos, etc.).
Incluye entre 5 y 15 filas de datos representativos.
Responde ÚNICAMENTE con el JSON válido, sin explicaciones ni bloques de código markdown.""",

    "pdf": """Eres el Agente DA de Arauco. Genera un informe ejecutivo en formato JSON estructurado.

REGLA ABSOLUTA: responde ÚNICAMENTE con JSON válido. Sin texto previo ni posterior. Sin bloques markdown.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
ESQUEMA JSON OBLIGATORIO
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
{
  "titulo": "Título principal del informe",
  "subtitulo": "Subtítulo o descripción breve",
  "fecha": "15 de abril de 2025",
  "area": "Subgerencia de Mejora Continua",
  "fuente": "Archivo o sistema de origen",
  "kpis": [
    {"label": "Nombre KPI", "valor": "1.234", "unidad": "unidad"}
  ],
  "resumen": "Texto del resumen ejecutivo. Puede tener varios párrafos separados por \\n\\n.",
  "secciones": [
    {
      "titulo": "Título de la sección",
      "tipo": "texto",
      "contenido": "Texto explicativo de la sección."
    },
    {
      "titulo": "Título de tabla",
      "tipo": "tabla",
      "encabezados": ["Col A", "Col B", "Col C"],
      "filas": [
        ["Valor 1", "Valor 2", "Valor 3"]
      ]
    }
  ],
  "conclusiones": "Texto de conclusiones y próximos pasos recomendados."
}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
REGLAS DE CONTENIDO
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
- kpis: entre 3 y 6 métricas clave derivadas del análisis
- secciones: entre 2 y 5 secciones; alterna texto y tablas según corresponda
- Las tablas muestran máximo 20 filas (top-20 por relevancia)
- EXCEL → usa stats para KPIs y totales; muestra_top20 para tablas
- PDF/WORD → extrae secciones, cifras y tablas del texto recibido
- NUNCA inventes cifras; usa formato numérico chileno: 1.234,5""",

    "gantt": """Eres el Agente DA de Arauco. Extrae las tareas del contexto recibido y devuelve un JSON estructurado.

REGLA ABSOLUTA: responde ÚNICAMENTE con JSON válido. Sin texto previo ni posterior. Sin bloques markdown.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
ESQUEMA JSON OBLIGATORIO
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
{
  "titulo": "Nombre del proyecto",
  "subtitulo": "Descripción breve del alcance",
  "fecha": "Rango de fechas legible, ej: Enero - Marzo 2025",
  "tareas": [
    {
      "id": "t1",
      "nombre": "Nombre de la tarea",
      "area": "EO",
      "responsable": "Nombre o equipo",
      "inicio": "YYYY-MM-DD",
      "fin": "YYYY-MM-DD",
      "avance": 75,
      "deps": ""
    }
  ]
}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
REGLAS
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
- area: EO (Excelencia Operacional) | TD (Transformación Digital) | IA (Inteligencia Artificial) | Gestión | Riesgo
- avance: número entero 0-100 (% completado)
- deps: id de tarea predecesora (string "t1") o "" si no tiene dependencia
- inicio y fin: formato YYYY-MM-DD obligatorio; fin debe ser posterior a inicio
- Mínimo 5 tareas, máximo 25
- Si el contexto tiene fechas y tareas reales: úsalas directamente
- Si no hay fechas: inferir un proyecto realista según el contexto recibido
- NUNCA inventes cifras de avance si hay datos reales""",

    "email": """Eres el Agente DA de Arauco. Redacta un correo profesional en formato JSON.

REGLA ABSOLUTA: responde ÚNICAMENTE con JSON válido. Sin texto previo ni posterior.

Esquema:
{
  "para": "destinatario@arauco.com",
  "cc": "",
  "asunto": "Asunto del correo",
  "cuerpo": "Cuerpo del correo en texto plano. Puede tener saltos de línea con \\n."
}

REGLAS:
- Tono profesional, conciso, en español
- Si el usuario indica destinatario, úsalo; si no, deja para="destinatario@arauco.com" como placeholder
- El cuerpo debe incluir saludo, desarrollo y cierre. La firma final siempre debe ser exactamente: "Saludos,\nLuciano"
- Si hay datos del análisis previo, resúmelos en el cuerpo de forma ejecutiva
- Nunca inventes cifras; usa solo las del contexto recibido""",

    "pptx": """Eres el Agente DA (Analista de Datos) de Arauco — Subgerencia de Mejora Continua.
Genera una presentación PowerPoint estructurada, profesional y lista para usar.

Responde ÚNICAMENTE con un objeto JSON válido (sin explicaciones, sin bloques de código).

Estructura JSON exacta:
{
  "titulo": "Título principal de la presentación",
  "subtitulo": "Subtítulo o descripción breve",
  "area": "Mejora Continua",
  "fecha": "YYYY-MM-DD",
  "autor": "Subgerencia de Mejora Continua — Arauco",
  "diapositivas": [
    {
      "tipo": "portada",
      "titulo": "Título de la presentación",
      "subtitulo": "Subtítulo",
      "nota": ""
    },
    {
      "tipo": "contenido",
      "titulo": "Título de la diapositiva",
      "bullets": [
        "Punto principal 1",
        "Punto principal 2",
        "Punto principal 3"
      ],
      "nota": "Notas del presentador opcionales"
    },
    {
      "tipo": "tabla",
      "titulo": "Título con tabla de datos",
      "encabezados": ["Col A", "Col B", "Col C"],
      "filas": [
        ["valor1", "valor2", "valor3"],
        ["valor4", "valor5", "valor6"]
      ],
      "nota": ""
    },
    {
      "tipo": "cierre",
      "titulo": "Conclusiones",
      "bullets": ["Conclusión 1", "Conclusión 2"],
      "nota": ""
    }
  ]
}

TIPOS DE DIAPOSITIVA:
- "portada": slide de apertura (solo titulo y subtitulo)
- "contenido": título + lista de bullets (máx 6 bullets por slide)
- "tabla": título + tabla de datos (encabezados + filas)
- "cierre": slide de cierre con conclusiones o próximos pasos

REGLAS:
- Mínimo 5 diapositivas, máximo 15
- La primera debe ser tipo "portada" y la última tipo "cierre"
- Bullets concisos, máx 12 palabras cada uno
- Si hay datos reales en el contexto, úsalos LITERALMENTE en las tablas
- NUNCA inventes KPIs ni cifras operacionales
- Fecha en formato YYYY-MM-DD""",
}


def extract_json(raw: str) -> dict:
    """
    Extrae y parsea el JSON de la respuesta de Claude de forma robusta.
    Maneja: bloques ```json ... ```, bloques ``` ... ```, y JSON directo.
    """
    text = raw.strip()
    # Quitar bloque de código si lo hay
    if text.startswith("```"):
        lines = text.split("\n")
        # quitar primera línea (```json o ```) y última línea (```)
        inner = "\n".join(lines[1:])
        if inner.rstrip().endswith("```"):
            inner = inner.rstrip()[:-3].rstrip()
        text = inner.strip()
    # Buscar el primer { y el último } para extraer solo el JSON
    start = text.find("{")
    end   = text.rfind("}")
    if start != -1 and end != -1 and end > start:
        text = text[start:end+1]
    try:
        return json.loads(text)
    except json.JSONDecodeError as exc:
        raise ValueError(f"JSON inválido en respuesta del modelo: {exc}") from exc


def store_html(html: str) -> str:
    """Guarda HTML en memoria y retorna la URL pública."""
    gid = uuid.uuid4().hex
    with _html_lock:
        _HTML_STORE[gid] = html
        if len(_HTML_STORE) > _MAX_STORE:
            _HTML_STORE.popitem(last=False)   # elimina el más antiguo
    return f"{PUBLIC_BASE}/g/{gid}"


def build_gantt(data: dict) -> str:
    """
    Genera un Gantt HTML/CSS puro — sin JavaScript ni CDN externos.
    Python calcula todas las posiciones. Funciona en cualquier browser,
    incluido el visor embebido de Telegram en iOS/Android.
    """
    from datetime import datetime, timedelta

    AREA_COLORS = {
        "EO":      "#BFB800",
        "TD":      "#EA7600",
        "IA":      "#2D6A9F",
        "Gestión": "#696158",
        "Riesgo":  "#C00000",
    }
    AREA_BG = {      # versión semitransparente para borde/fondo
        "EO":      "#f0edbb",
        "TD":      "#fde8cc",
        "IA":      "#cde0f5",
        "Gestión": "#e0dcd8",
        "Riesgo":  "#f5cccc",
    }

    titulo    = data.get("titulo", "Carta Gantt")
    subtitulo = data.get("subtitulo", "Arauco — Mejora Continua")
    fecha     = data.get("fecha", "")
    tareas    = data.get("tareas", [])

    # ── Parsear y validar tareas ─────────────────────────────
    parsed = []
    for t in tareas:
        try:
            ini = datetime.strptime(t["inicio"], "%Y-%m-%d").date()
            fin = datetime.strptime(t["fin"],    "%Y-%m-%d").date()
            if fin <= ini:
                fin = ini + timedelta(days=1)
        except (KeyError, ValueError):
            continue
        try:
            avance = max(0, min(100, int(float(t.get("avance", 0) or 0))))
        except (TypeError, ValueError):
            avance = 0
        parsed.append({
            "nombre":      t.get("nombre", "Tarea"),
            "area":        t.get("area", "Gestión"),
            "responsable": t.get("responsable", ""),
            "ini": ini, "fin": fin, "avance": avance,
        })

    if not parsed:
        parsed = [{"nombre": "Sin tareas", "area": "Gestión",
                   "responsable": "", "ini": datetime.today().date(),
                   "fin": datetime.today().date() + timedelta(days=7), "avance": 0}]

    # ── Rango total del proyecto ─────────────────────────────
    total_ini = min(t["ini"] for t in parsed)
    total_fin = max(t["fin"] for t in parsed)
    total_days = max((total_fin - total_ini).days, 1)

    # ── Semanas del eje X ────────────────────────────────────
    # Empezar en lunes anterior al total_ini
    from datetime import date
    axis_start = total_ini - timedelta(days=total_ini.weekday())
    axis_end   = total_fin + timedelta(days=(6 - total_fin.weekday()))
    axis_days  = max((axis_end - axis_start).days, 1)

    semanas = []
    cur = axis_start
    while cur < axis_end:
        semanas.append(cur)
        cur += timedelta(days=7)

    def pct(d):
        """% de posición de una fecha dentro del eje."""
        return round((d - axis_start).days / axis_days * 100, 3)

    # ── Avance promedio ──────────────────────────────────────
    prom_avance = round(sum(t["avance"] for t in parsed) / len(parsed))

    # ── Leyenda HTML ─────────────────────────────────────────
    leyenda_html = "".join(
        f'<span class="leg"><span class="leg-dot" style="background:{c}"></span>{a}</span>'
        for a, c in AREA_COLORS.items()
    )

    # ── Marcas de semana en el eje ───────────────────────────
    MESES = ["Ene","Feb","Mar","Abr","May","Jun","Jul","Ago","Sep","Oct","Nov","Dic"]
    axis_ticks = ""
    for s in semanas:
        left = pct(s)
        lbl  = f"{s.day} {MESES[s.month-1]}"
        axis_ticks += (
            f'<div class="tick" style="left:{left}%">'
            f'<span class="tick-lbl">{lbl}</span></div>'
        )

    # ── Línea "Hoy" ──────────────────────────────────────────
    hoy = datetime.today().date()
    hoy_html = ""
    if axis_start <= hoy <= axis_end:
        hoy_html = (
            f'<div class="hoy-line" style="left:{pct(hoy)}%" title="Hoy">'
            f'<span class="hoy-lbl">Hoy</span></div>'
        )

    # ── Filas de tareas ──────────────────────────────────────
    filas_html = ""
    for i, t in enumerate(parsed):
        color  = AREA_COLORS.get(t["area"], "#696158")
        bg     = AREA_BG.get(t["area"], "#e8e8e8")
        row_bg = "#ffffff" if i % 2 == 0 else "#f7f6f4"
        left_  = pct(t["ini"])
        width_ = max(pct(t["fin"]) - left_, 0.5)
        avance_w = round(width_ * t["avance"] / 100, 3)
        dur_dias = (t["fin"] - t["ini"]).days
        ini_str  = t["ini"].strftime("%-d %b")
        fin_str  = t["fin"].strftime("%-d %b")
        resp_str = f" · {t['responsable']}" if t["responsable"] else ""

        filas_html += f"""
<div class="row" style="background:{row_bg}">
  <div class="row-label">
    <span class="dot" style="background:{color}"></span>
    <div>
      <div class="task-name">{t['nombre']}</div>
      <div class="task-meta">{t['area']}{resp_str}</div>
    </div>
  </div>
  <div class="row-bar">
    <div class="bar-bg" style="left:{left_}%;width:{width_}%;background:{bg};border-color:{color}">
      <div class="bar-fg" style="width:{avance_w / width_ * 100 if width_ else 0}%;background:{color}"></div>
      <span class="bar-lbl">{ini_str} → {fin_str} &nbsp;|&nbsp; {t['avance']}% &nbsp;({dur_dias}d)</span>
    </div>
    {hoy_html}
  </div>
</div>"""

    html = f"""<!DOCTYPE html>
<html lang="es">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0">
<title>{titulo}</title>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:Arial,Helvetica,sans-serif;background:#f0ede9;color:#333}}
.header{{background:#696158;color:#fff;padding:16px 20px}}
.header h1{{font-size:1.1rem;font-weight:900}}
.header p{{font-size:0.75rem;opacity:.8;margin-top:3px}}
.controls{{background:#fff;padding:10px 20px;border-bottom:1px solid #EDEAE6;
           display:flex;flex-wrap:wrap;align-items:center;gap:12px}}
.badge{{background:#BFB800;color:#fff;border-radius:12px;padding:3px 12px;
        font-size:0.78rem;font-weight:700;white-space:nowrap}}
.leg{{display:inline-flex;align-items:center;gap:4px;font-size:0.72rem;color:#555}}
.leg-dot{{width:10px;height:10px;border-radius:2px;flex-shrink:0}}
.gantt{{overflow-x:auto;padding:12px 20px}}
.axis{{position:relative;height:28px;margin-left:160px;border-bottom:2px solid #ccc;margin-bottom:0}}
.tick{{position:absolute;top:0;transform:translateX(-50%);display:flex;flex-direction:column;align-items:center}}
.tick::before{{content:'';display:block;width:1px;height:8px;background:#ccc}}
.tick-lbl{{font-size:0.62rem;color:#888;white-space:nowrap;margin-top:2px}}
.row{{display:flex;align-items:center;min-height:44px;border-bottom:1px solid #EDEAE6}}
.row-label{{width:160px;min-width:160px;padding:6px 10px 6px 4px;display:flex;
            align-items:flex-start;gap:6px;flex-shrink:0}}
.dot{{width:8px;height:8px;border-radius:2px;margin-top:4px;flex-shrink:0}}
.task-name{{font-size:0.75rem;font-weight:700;line-height:1.3;color:#333}}
.task-meta{{font-size:0.65rem;color:#888;margin-top:1px}}
.row-bar{{flex:1;position:relative;height:44px;min-width:0}}
.bar-bg{{position:absolute;top:50%;transform:translateY(-50%);height:22px;
         border-radius:4px;border:1.5px solid;overflow:hidden;min-width:4px}}
.bar-fg{{height:100%;border-radius:3px;opacity:0.9}}
.bar-lbl{{position:absolute;left:4px;top:50%;transform:translateY(-50%);
          font-size:0.6rem;color:#333;white-space:nowrap;font-weight:600;
          pointer-events:none;mix-blend-mode:multiply}}
.hoy-line{{position:absolute;top:0;bottom:0;width:2px;background:#C00000;
           opacity:.7;pointer-events:none;z-index:10}}
.hoy-lbl{{position:absolute;top:2px;left:3px;font-size:0.58rem;
          color:#C00000;white-space:nowrap;font-weight:700}}
.footer{{background:#696158;color:#fff;text-align:center;
         padding:12px;font-size:0.72rem;opacity:.9;margin-top:16px}}
</style>
</head>
<body>
<div class="header" style="display:flex;align-items:center;gap:14px">
  <img src="{_ARAUCO_LOGO_URL}" alt="Arauco" height="28" style="flex-shrink:0;object-fit:contain">
  <div><h1>{titulo}</h1><p>{subtitulo} &nbsp;·&nbsp; {fecha}</p></div>
</div>
<div class="controls">
  <span class="badge">Avance promedio: {prom_avance}%</span>
  {leyenda_html}
</div>
<div class="gantt">
  <div class="axis">{axis_ticks}</div>
  {filas_html}
</div>
<div class="footer">Arauco — Subgerencia de Mejora Continua &nbsp;|&nbsp; {fecha}</div>
</body>
</html>"""

    url = store_html(html)
    return url


def build_pdf(data: dict) -> io.BytesIO:
    """Genera un PDF ejecutivo Arauco a partir del JSON estructurado."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import mm
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        HRFlowable, KeepTogether
    )
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY

    GRIS    = colors.HexColor("#696158")
    VERDE   = colors.HexColor("#BFB800")
    NARANJA = colors.HexColor("#EA7600")
    CREMA   = colors.HexColor("#EDEAE6")
    BLANCO  = colors.white
    NEGRO   = colors.HexColor("#222222")
    GRIS_L  = colors.HexColor("#999999")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=20*mm, rightMargin=20*mm,
        topMargin=18*mm, bottomMargin=18*mm,
        title=data.get("titulo", "Informe Arauco"),
    )

    # Estilos
    s_titulo  = ParagraphStyle("titulo",  fontName="Helvetica-Bold",   fontSize=18, textColor=GRIS,  leading=26, spaceAfter=4*mm)
    s_sub     = ParagraphStyle("sub",     fontName="Helvetica",        fontSize=10, textColor=GRIS_L, leading=15, spaceAfter=8*mm)
    s_meta    = ParagraphStyle("meta",    fontName="Helvetica",        fontSize=8,  textColor=GRIS_L, spaceAfter=4*mm)
    s_seccion = ParagraphStyle("seccion", fontName="Helvetica-Bold",   fontSize=12, textColor=GRIS,  leading=18, spaceBefore=8*mm, spaceAfter=4*mm, borderPad=2, leftIndent=4*mm)
    s_body    = ParagraphStyle("body",    fontName="Helvetica",        fontSize=9,  textColor=NEGRO, leading=14, alignment=TA_JUSTIFY, spaceAfter=5*mm)
    s_concl   = ParagraphStyle("concl",   fontName="Helvetica-Oblique",fontSize=9,  textColor=NEGRO, leading=14, alignment=TA_JUSTIFY, spaceAfter=5*mm)
    s_footer  = ParagraphStyle("footer",  fontName="Helvetica",        fontSize=7,  textColor=GRIS_L, alignment=TA_CENTER)
    s_kpi_val = ParagraphStyle("kpi_val", fontName="Helvetica-Bold",   fontSize=16, textColor=GRIS,  alignment=TA_CENTER, leading=18)
    s_kpi_lbl = ParagraphStyle("kpi_lbl", fontName="Helvetica",        fontSize=7,  textColor=GRIS_L, alignment=TA_CENTER, spaceAfter=0)
    s_kpi_uni = ParagraphStyle("kpi_uni", fontName="Helvetica",        fontSize=7,  textColor=GRIS_L, alignment=TA_CENTER)

    from reportlab.platypus import Image as RLImage

    story = []

    # ── BANDA LOGO ───────────────────────────────────────────
    logo_bytes = _get_logo_bytes()
    s_arauco_txt = ParagraphStyle("arauco_txt", fontName="Helvetica-Bold", fontSize=14,
                                  textColor=BLANCO, alignment=TA_LEFT)
    s_area_txt   = ParagraphStyle("area_txt",   fontName="Helvetica",      fontSize=8,
                                  textColor=BLANCO, alignment=TA_RIGHT)
    page_w = A4[0] - 40*mm
    if logo_bytes:
        logo_cell = RLImage(io.BytesIO(logo_bytes), width=30*mm, height=10*mm)
    else:
        logo_cell = Paragraph("ARAUCO", s_arauco_txt)
    area_cell = Paragraph(data.get("area", "Subgerencia de Mejora Continua"), s_area_txt)
    logo_band = Table([[logo_cell, area_cell]], colWidths=[page_w * 0.5, page_w * 0.5])
    logo_band.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, -1), GRIS),
        ("TOPPADDING",    (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ("LEFTPADDING",   (0, 0), (0,  -1), 8),
        ("RIGHTPADDING",  (-1, 0), (-1, -1), 8),
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
    ]))
    story.append(logo_band)
    story.append(Spacer(1, 4*mm))

    # ── HEADER ──────────────────────────────────────────────
    story.append(Paragraph(data.get("titulo", "Informe Ejecutivo"), s_titulo))
    if data.get("subtitulo"):
        story.append(Paragraph(data["subtitulo"], s_sub))
    story.append(Paragraph(
        f"<b>Área:</b> {data.get('area','Mejora Continua')} &nbsp;|&nbsp; "
        f"<b>Fecha:</b> {data.get('fecha','')} &nbsp;|&nbsp; "
        f"<b>Fuente:</b> {data.get('fuente','')}",
        s_meta
    ))
    story.append(HRFlowable(width="100%", thickness=2, color=GRIS, spaceAfter=5*mm))

    # ── KPIs ────────────────────────────────────────────────
    kpis = data.get("kpis", [])
    if kpis:
        n = len(kpis)
        col_w = (A4[0] - 40*mm) / n
        kpi_data = [[
            Paragraph(k.get("valor", ""), s_kpi_val) for k in kpis
        ], [
            Paragraph(k.get("label", ""), s_kpi_lbl) for k in kpis
        ], [
            Paragraph(k.get("unidad", ""), s_kpi_uni) for k in kpis
        ]]
        kpi_table = Table(kpi_data, colWidths=[col_w]*n, rowHeights=[20*mm, 6*mm, 5*mm])
        kpi_table.setStyle(TableStyle([
            ("BOX",         (0,0), (-1,-1), 0.5, CREMA),
            ("INNERGRID",   (0,0), (-1,-1), 0.5, CREMA),
            ("BACKGROUND",  (0,0), (-1,-1), colors.white),
            ("TOPPADDING",  (0,0), (-1,-1), 4),
            ("BOTTOMPADDING",(0,0),(-1,-1), 4),
            ("LINEBELOW",   (0,0), (-1,0),  1.5, VERDE),
        ]))
        story.append(kpi_table)
        story.append(Spacer(1, 5*mm))

    # ── RESUMEN ─────────────────────────────────────────────
    resumen = data.get("resumen", "")
    if resumen:
        story.append(Paragraph("Resumen Ejecutivo", s_seccion))
        story.append(HRFlowable(width="100%", thickness=1, color=VERDE, spaceAfter=3*mm))
        s_resumen = ParagraphStyle("resumen", fontName="Helvetica", fontSize=9,
                                   textColor=NEGRO, leading=14, alignment=TA_JUSTIFY,
                                   spaceAfter=4*mm)
        page_w = A4[0] - 40*mm
        parrafos = [Paragraph(p.strip(), s_resumen)
                    for p in resumen.split("\n\n") if p.strip()]
        resumen_inner = [[parrafos]]
        resumen_table = Table(resumen_inner, colWidths=[page_w])
        resumen_table.setStyle(TableStyle([
            ("BACKGROUND",    (0,0), (-1,-1), CREMA),
            ("BOX",           (0,0), (-1,-1), 0,    CREMA),
            ("LEFTPADDING",   (0,0), (-1,-1), 8),
            ("RIGHTPADDING",  (0,0), (-1,-1), 8),
            ("TOPPADDING",    (0,0), (-1,-1), 6),
            ("BOTTOMPADDING", (0,0), (-1,-1), 6),
            ("LINEBELOW",     (0,0), (-1,0),  2, VERDE),
        ]))
        story.append(resumen_table)
        story.append(Spacer(1, 4*mm))

    # ── SECCIONES ───────────────────────────────────────────
    for sec in data.get("secciones", []):
        titulo_sec = sec.get("titulo", "")
        tipo       = sec.get("tipo", "texto")

        bloque = [
            Paragraph(titulo_sec, s_seccion),
            HRFlowable(width="100%", thickness=1, color=VERDE, spaceAfter=3*mm),
        ]

        if tipo == "texto":
            contenido = sec.get("contenido", "")
            for parr in contenido.split("\n\n"):
                parr = parr.strip()
                if parr:
                    bloque.append(Paragraph(parr, s_body))

        elif tipo == "tabla":
            encab = sec.get("encabezados", [])
            filas = sec.get("filas", [])
            if encab:
                page_w = A4[0] - 40*mm
                col_w  = page_w / len(encab)
                t_data = [[Paragraph(str(h), ParagraphStyle("th", fontName="Helvetica-Bold",
                           fontSize=8, textColor=BLANCO)) for h in encab]]
                for fila in filas:
                    t_data.append([Paragraph(str(v), ParagraphStyle("td", fontName="Helvetica",
                                   fontSize=8, textColor=NEGRO)) for v in fila])
                t = Table(t_data, colWidths=[col_w]*len(encab), repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND",    (0,0),  (-1,0),  GRIS),
                    ("TEXTCOLOR",     (0,0),  (-1,0),  BLANCO),
                    ("ROWBACKGROUNDS",(0,1),  (-1,-1), [colors.white, CREMA]),
                    ("GRID",          (0,0),  (-1,-1), 0.3, colors.HexColor("#dddddd")),
                    ("TOPPADDING",    (0,0),  (-1,-1), 4),
                    ("BOTTOMPADDING", (0,0),  (-1,-1), 4),
                    ("LEFTPADDING",   (0,0),  (-1,-1), 6),
                ]))
                bloque.append(t)

        story.append(KeepTogether(bloque))
        story.append(Spacer(1, 3*mm))

    # ── CONCLUSIONES ────────────────────────────────────────
    conclusiones = data.get("conclusiones", "")
    if conclusiones:
        bloque_c = [
            Paragraph("Conclusiones y Recomendaciones", s_seccion),
            HRFlowable(width="100%", thickness=1, color=NARANJA, spaceAfter=3*mm),
        ]
        for parr in conclusiones.split("\n\n"):
            parr = parr.strip()
            if parr:
                bloque_c.append(Paragraph(parr, s_concl))
        story.append(KeepTogether(bloque_c))

    # ── FOOTER ──────────────────────────────────────────────
    story.append(Spacer(1, 8*mm))
    story.append(HRFlowable(width="100%", thickness=0.5, color=CREMA, spaceAfter=3*mm))
    story.append(Paragraph(
        f"Arauco — Subgerencia de Mejora Continua &nbsp;|&nbsp; {data.get('fecha','')} &nbsp;|&nbsp; {data.get('fuente','')}",
        s_footer
    ))

    doc.build(story)
    buf.seek(0)
    return buf


def build_pdf_imagenes(data: dict, images: list) -> io.BytesIO:
    """PDF con imágenes reales: portada → 1 página por imagen → conclusiones."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import mm
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, HRFlowable,
        Image as RLImage, PageBreak, KeepTogether
    )
    from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER
    from reportlab.lib.utils import ImageReader

    GRIS   = colors.HexColor("#696158")
    VERDE  = colors.HexColor("#BFB800")
    NEGRO  = colors.HexColor("#222222")
    GRIS_L = colors.HexColor("#999999")
    CREMA  = colors.HexColor("#EDEAE6")

    PAGE_W, PAGE_H = A4
    MARGIN = 20 * mm
    USABLE_W = PAGE_W - 2 * MARGIN

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=MARGIN, rightMargin=MARGIN,
                            topMargin=18*mm, bottomMargin=18*mm,
                            title=data.get("titulo", "Informe Visual Arauco"))

    s_titulo  = ParagraphStyle("t",  fontName="Helvetica-Bold",   fontSize=20,
                                textColor=GRIS, leading=26, spaceAfter=4*mm)
    s_sub     = ParagraphStyle("s",  fontName="Helvetica",        fontSize=11,
                                textColor=GRIS_L, spaceAfter=6*mm)
    s_meta    = ParagraphStyle("m",  fontName="Helvetica",        fontSize=8,
                                textColor=GRIS_L, spaceAfter=4*mm)
    s_seccion = ParagraphStyle("sc", fontName="Helvetica-Bold",   fontSize=12,
                                textColor=GRIS, spaceBefore=4*mm, spaceAfter=2*mm)
    s_body    = ParagraphStyle("b",  fontName="Helvetica",        fontSize=9,
                                textColor=NEGRO, leading=14,
                                alignment=TA_JUSTIFY, spaceAfter=3*mm)
    s_bullet  = ParagraphStyle("bu", fontName="Helvetica",        fontSize=9,
                                textColor=NEGRO, leading=13,
                                leftIndent=8*mm, spaceAfter=2*mm)
    s_concl   = ParagraphStyle("co", fontName="Helvetica-Oblique",fontSize=9,
                                textColor=NEGRO, leading=14, spaceAfter=3*mm)
    s_caption = ParagraphStyle("cp", fontName="Helvetica",        fontSize=7,
                                textColor=GRIS_L, alignment=TA_CENTER, spaceAfter=3*mm)

    story = []

    # ── PORTADA ──────────────────────────────────────────────────────────
    story.append(Paragraph(data.get("titulo", "Registro Visual"), s_titulo))
    if data.get("subtitulo"):
        story.append(Paragraph(data["subtitulo"], s_sub))
    story.append(Paragraph(
        f"<b>Área:</b> {data.get('area', 'Mejora Continua')} &nbsp;|&nbsp; "
        f"<b>Fecha:</b> {data.get('fecha', '')} &nbsp;|&nbsp; "
        f"<b>Imágenes:</b> {len(images)}",
        s_meta
    ))
    story.append(HRFlowable(width="100%", thickness=2, color=GRIS, spaceAfter=6*mm))

    # Índice rápido
    for i, img_meta in enumerate(data.get("imagenes", [])[:len(images)]):
        story.append(Paragraph(f"{i+1}. {img_meta.get('titulo', f'Imagen {i+1}')}", s_body))
    story.append(PageBreak())

    # ── UNA PÁGINA POR IMAGEN ────────────────────────────────────────────
    img_metas = data.get("imagenes", [])
    for i, img in enumerate(images):
        meta = img_metas[i] if i < len(img_metas) else {}

        # Calcular dimensiones manteniendo proporción
        try:
            raw = base64.b64decode(img["b64"])
            ir  = ImageReader(io.BytesIO(raw))
            px_w, px_h = ir.getSize()
            aspect   = px_h / px_w
            max_img_h = 130 * mm
            disp_w   = USABLE_W
            disp_h   = min(disp_w * aspect, max_img_h)
            if disp_h == max_img_h:
                disp_w = disp_h / aspect
            img_flow = RLImage(io.BytesIO(raw), width=disp_w, height=disp_h)
        except Exception:
            img_flow = None

        bloque = []
        bloque.append(Paragraph(meta.get("titulo", f"Imagen {i+1}"), s_seccion))
        bloque.append(HRFlowable(width="100%", thickness=1, color=VERDE, spaceAfter=3*mm))
        if img_flow:
            bloque.append(img_flow)
            bloque.append(Paragraph(f"Imagen {i+1} de {len(images)}", s_caption))
        if meta.get("descripcion"):
            bloque.append(Paragraph(meta["descripcion"], s_body))
        for h in meta.get("hallazgos", []):
            bloque.append(Paragraph(f"• {h}", s_bullet))

        story.append(KeepTogether(bloque[:3]))   # título + línea + imagen juntos
        story.extend(bloque[3:])
        if i < len(images) - 1:
            story.append(PageBreak())

    # ── CONCLUSIONES ─────────────────────────────────────────────────────
    conclusiones = data.get("conclusiones", [])
    if conclusiones:
        story.append(PageBreak())
        story.append(Paragraph("Conclusiones y Recomendaciones", s_seccion))
        story.append(HRFlowable(width="100%", thickness=1, color=GRIS, spaceAfter=4*mm))
        for c in conclusiones:
            story.append(Paragraph(f"• {c}", s_concl))

    doc.build(story)
    buf.seek(0)
    return buf


def build_notas_pdf(notas: list, titulo: str = "Notas") -> io.BytesIO:
    """PDF directo desde lista de notas — sin llamada a Claude."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import mm
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    from reportlab.lib.enums import TA_JUSTIFY

    GRIS   = colors.HexColor("#696158")
    NEGRO  = colors.HexColor("#222222")
    GRIS_L = colors.HexColor("#999999")
    SEP    = colors.HexColor("#DDDDDD")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=20*mm, rightMargin=20*mm,
                            topMargin=18*mm, bottomMargin=18*mm,
                            title=titulo)

    s_titulo = ParagraphStyle("titulo", fontName="Helvetica-Bold", fontSize=16,
                               textColor=GRIS, leading=22, spaceAfter=3*mm)
    s_meta   = ParagraphStyle("meta",   fontName="Helvetica", fontSize=8,
                               textColor=GRIS_L, spaceAfter=5*mm)
    s_ts     = ParagraphStyle("ts",     fontName="Helvetica", fontSize=8,
                               textColor=GRIS_L, spaceAfter=1*mm)
    s_body   = ParagraphStyle("body",   fontName="Helvetica", fontSize=9,
                               textColor=NEGRO, leading=14, alignment=TA_JUSTIFY,
                               spaceAfter=4*mm)

    story = []
    story.append(Paragraph(titulo, s_titulo))
    story.append(Paragraph(
        f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')} &nbsp;|&nbsp; "
        f"{len(notas)} nota(s) &nbsp;|&nbsp; Arauco — Mejora Continua",
        s_meta
    ))
    story.append(HRFlowable(width="100%", thickness=2, color=GRIS, spaceAfter=5*mm))

    for n in notas:
        story.append(Paragraph(f"Nota {n['n']} — {n['fecha']}", s_ts))
        body_text = n['texto'].replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        story.append(Paragraph(body_text, s_body))
        story.append(HRFlowable(width="100%", thickness=0.5, color=SEP, spaceAfter=3*mm))

    doc.build(story)
    buf.seek(0)
    return buf


def build_notas_docx(notas: list, titulo: str = "Notas") -> io.BytesIO:
    """Word (.docx) desde lista de notas — formato compatible con importación en OneNote."""
    doc = DocxDocument()

    doc.add_heading(titulo, 0)
    doc.add_paragraph(
        f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}  |  "
        f"{len(notas)} nota(s)  |  Arauco — Mejora Continua"
    )
    doc.add_paragraph("")

    for n in notas:
        doc.add_heading(f"Nota {n['n']} — {n['fecha']}", level=2)
        doc.add_paragraph(n['texto'])
        doc.add_paragraph("")

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def build_excel(data: dict) -> io.BytesIO:
    """Crea un archivo Excel a partir del JSON generado por Claude."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = data.get("titulo", "Datos")[:31]

    # Encabezados con estilo
    from openpyxl.styles import Font, PatternFill, Alignment
    header_fill = PatternFill(start_color="2D6A4F", end_color="2D6A4F", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)

    for col, header in enumerate(data["encabezados"], start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")
        ws.column_dimensions[cell.column_letter].width = max(len(str(header)) + 4, 12)

    for row_idx, fila in enumerate(data["filas"], start=2):
        for col_idx, valor in enumerate(fila, start=1):
            ws.cell(row=row_idx, column=col_idx, value=valor)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf



def build_pptx(data: dict) -> io.BytesIO:
    """Genera un archivo PowerPoint con diseño Arauco a partir del JSON estructurado."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    GRIS    = RGBColor(0x69, 0x61, 0x58)
    AMARILLO = RGBColor(0xBF, 0xB8, 0x00)
    NARANJA = RGBColor(0xEA, 0x76, 0x00)
    BLANCO  = RGBColor(0xFF, 0xFF, 0xFF)
    NEGRO   = RGBColor(0x22, 0x22, 0x22)
    GRIS_L  = RGBColor(0x99, 0x99, 0x99)
    CREMA   = RGBColor(0xED, 0xEA, 0xE6)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BLANK = prs.slide_layouts[6]  # layout en blanco

    def _set_bg(slide, color: RGBColor):
        from pptx.oxml.ns import qn
        from lxml import etree
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text_box(slide, text, left, top, width, height,
                      font_size=18, bold=False, color=NEGRO,
                      align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txBox

    def _add_rect(slide, left, top, width, height, fill_color: RGBColor):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    for slide_data in data.get("diapositivas", []):
        tipo  = slide_data.get("tipo", "contenido")
        slide = prs.slides.add_slide(BLANK)

        # ── PORTADA ───────────────────────────────────────────────────
        if tipo == "portada":
            _set_bg(slide, GRIS)
            _add_rect(slide, 0, 0, 13.33, 0.12, AMARILLO)    # línea top
            _add_rect(slide, 0, 7.38, 13.33, 0.12, AMARILLO) # línea bottom
            logo_bytes = _get_logo_bytes()
            if logo_bytes:
                slide.shapes.add_picture(
                    io.BytesIO(logo_bytes),
                    Inches(1.0), Inches(0.3), width=Inches(2.2)
                )
            else:
                _add_text_box(slide, "ARAUCO", 1.0, 0.3, 4, 0.7,
                              font_size=13, bold=True, color=AMARILLO)
            _add_text_box(slide, slide_data.get("titulo", ""),
                          1.0, 2.0, 11.33, 2.0,
                          font_size=36, bold=True, color=BLANCO, wrap=True)
            _add_text_box(slide, slide_data.get("subtitulo", ""),
                          1.0, 4.2, 11.33, 1.2,
                          font_size=20, bold=False, color=CREMA, wrap=True)
            _add_text_box(slide,
                          f"{data.get('area','')}  |  {data.get('fecha','')}  |  {data.get('autor','')}",
                          1.0, 6.6, 11.33, 0.6,
                          font_size=10, bold=False, color=GRIS_L)

        # ── CIERRE ────────────────────────────────────────────────────
        elif tipo == "cierre":
            _set_bg(slide, GRIS)
            _add_rect(slide, 0, 0, 13.33, 0.12, AMARILLO)
            _add_rect(slide, 0, 7.38, 13.33, 0.12, AMARILLO)
            _add_text_box(slide, slide_data.get("titulo", "Conclusiones"),
                          1.0, 1.5, 11.33, 1.2,
                          font_size=30, bold=True, color=BLANCO)
            bullets = slide_data.get("bullets", [])
            y = 3.0
            for bullet in bullets:
                _add_text_box(slide, f"• {bullet}", 1.2, y, 10.8, 0.55,
                              font_size=16, color=CREMA, wrap=True)
                y += 0.6

        # ── TABLA ─────────────────────────────────────────────────────
        elif tipo == "tabla":
            _set_bg(slide, BLANCO)
            _add_rect(slide, 0, 0, 13.33, 1.1, GRIS)
            _add_text_box(slide, slide_data.get("titulo", ""),
                          0.3, 0.18, 12.73, 0.75,
                          font_size=22, bold=True, color=BLANCO)
            headers = slide_data.get("encabezados", [])
            rows    = slide_data.get("filas", [])
            if headers:
                cols = len(headers)
                rows_total = 1 + len(rows)
                tbl_left   = Inches(0.4)
                tbl_top    = Inches(1.3)
                tbl_width  = Inches(12.53)
                tbl_height = Inches(min(rows_total * 0.45, 5.5))
                table = slide.shapes.add_table(
                    rows_total, cols, tbl_left, tbl_top, tbl_width, tbl_height
                ).table
                col_w = int(tbl_width / cols)
                for c in range(cols):
                    table.columns[c].width = col_w

                # Encabezado
                for c, h in enumerate(headers):
                    cell = table.cell(0, c)
                    cell.text = h
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = GRIS
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0] if p.runs else p.add_run()
                    run.font.bold  = True
                    run.font.size  = Pt(11)
                    run.font.color.rgb = BLANCO

                # Filas de datos
                for r, fila in enumerate(rows, start=1):
                    bg = CREMA if r % 2 == 0 else BLANCO
                    for c, val in enumerate(fila[:cols]):
                        cell = table.cell(r, c)
                        cell.text = str(val)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = bg
                        p = cell.text_frame.paragraphs[0]
                        run = p.runs[0] if p.runs else p.add_run()
                        run.font.size  = Pt(10)
                        run.font.color.rgb = NEGRO

        # ── CONTENIDO (default) ───────────────────────────────────────
        else:
            _set_bg(slide, BLANCO)
            _add_rect(slide, 0, 0, 13.33, 1.1, GRIS)
            _add_text_box(slide, slide_data.get("titulo", ""),
                          0.3, 0.18, 12.73, 0.75,
                          font_size=22, bold=True, color=BLANCO)
            _add_rect(slide, 0.4, 1.15, 0.06, 5.9, AMARILLO)  # barra lateral
            bullets = slide_data.get("bullets", [])
            y = 1.3
            for bullet in bullets[:6]:
                _add_text_box(slide, f"  {bullet}", 0.6, y, 12.33, 0.65,
                              font_size=16, color=NEGRO, wrap=True)
                y += 0.72

        # Notas del presentador
        nota = slide_data.get("nota", "")
        if nota:
            slide.notes_slide.notes_text_frame.text = nota

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def build_pptx_imagenes(data: dict, images: list) -> io.BytesIO:
    """PPT con imágenes reales: portada → 1 slide por imagen → cierre."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    GRIS    = RGBColor(0x69, 0x61, 0x58)
    AMARILLO = RGBColor(0xBF, 0xB8, 0x00)
    BLANCO  = RGBColor(0xFF, 0xFF, 0xFF)
    NEGRO   = RGBColor(0x22, 0x22, 0x22)
    CREMA   = RGBColor(0xED, 0xEA, 0xE6)
    GRIS_L  = RGBColor(0x99, 0x99, 0x99)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def _bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _rect(slide, l, t, w, h, color):
        s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

    def _txt(slide, text, l, t, w, h, size=16, bold=False, color=NEGRO,
             align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]; p.alignment = align
        r  = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    # ── PORTADA ──────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    _bg(sl, GRIS)
    _rect(sl, 0, 0, 13.33, 0.12, AMARILLO)
    _rect(sl, 0, 7.38, 13.33, 0.12, AMARILLO)
    _txt(sl, "ARAUCO", 1.0, 0.4, 11, 0.6, size=13, bold=True, color=AMARILLO)
    _txt(sl, data.get("titulo", "Registro Visual"),
         1.0, 1.8, 11.33, 2.0, size=36, bold=True, color=BLANCO)
    _txt(sl, data.get("subtitulo", ""),
         1.0, 4.0, 11.33, 1.0, size=20, color=CREMA)
    _txt(sl, f"{data.get('area','')}  |  {data.get('fecha','')}  |  {data.get('autor','')}",
         1.0, 6.6, 11.33, 0.6, size=10, color=GRIS_L)

    # ── SLIDE POR IMAGEN ─────────────────────────────────────────────────
    img_metas = data.get("imagenes", [])
    for i, img in enumerate(images):
        meta  = img_metas[i] if i < len(img_metas) else {}
        sl    = prs.slides.add_slide(BLANK)
        _bg(sl, BLANCO)
        _rect(sl, 0, 0, 13.33, 1.1, GRIS)
        _txt(sl, meta.get("titulo", f"Imagen {i+1}"),
             0.3, 0.18, 12.73, 0.75, size=22, bold=True, color=BLANCO)
        _rect(sl, 0.4, 1.15, 0.06, 6.0, AMARILLO)  # barra lateral

        # Imagen izquierda
        try:
            img_bytes  = base64.b64decode(img["b64"])
            img_stream = io.BytesIO(img_bytes)
            sl.shapes.add_picture(img_stream, Inches(0.55), Inches(1.25),
                                  Inches(6.8), Inches(5.8))
        except Exception:
            _txt(sl, "[imagen no disponible]", 0.6, 3.5, 6.5, 0.5, size=12, color=GRIS_L)

        # Texto derecha
        desc = meta.get("descripcion", "")
        if desc:
            _txt(sl, desc, 7.6, 1.3, 5.4, 1.8, size=13, color=NEGRO)

        hallazgos = meta.get("hallazgos", [])
        y = 3.3
        for h in hallazgos[:5]:
            _txt(sl, f"• {h}", 7.6, y, 5.4, 0.6, size=12, color=NEGRO)
            y += 0.65

        nota = meta.get("nota", "")
        if nota:
            sl.notes_slide.notes_text_frame.text = nota

    # ── CIERRE ───────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(BLANK)
    _bg(sl, GRIS)
    _rect(sl, 0, 0, 13.33, 0.12, AMARILLO)
    _rect(sl, 0, 7.38, 13.33, 0.12, AMARILLO)
    _txt(sl, "Conclusiones", 1.0, 1.3, 11.33, 1.0, size=30, bold=True, color=BLANCO)
    y = 2.8
    for c in data.get("conclusiones", [])[:6]:
        _txt(sl, f"• {c}", 1.2, y, 10.8, 0.6, size=16, color=CREMA)
        y += 0.7

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def send_email_outlook(data: dict) -> None:
    """Envía un correo via SendGrid API."""
    import requests

    api_key     = os.environ.get("SENDGRID_API_KEY", "")
    sender_email = os.environ.get("SENDER_EMAIL", "")
    if not api_key or not sender_email:
        raise ValueError("Faltan variables de entorno SENDGRID_API_KEY y/o SENDER_EMAIL.")

    to_list = [{"email": data["para"]}]
    payload = {
        "personalizations": [{"to": to_list}],
        "from": {"email": sender_email},
        "subject": data["asunto"],
        "content": [{"type": "text/plain", "value": data["cuerpo"]}],
    }
    if data.get("cc"):
        payload["personalizations"][0]["cc"] = [{"email": data["cc"]}]
    if data.get("adjuntos"):
        payload["attachments"] = [
            {
                "content":     img["b64"],
                "filename":    f"imagen-{i + 1}.jpg",
                "type":        img.get("media_type", "image/jpeg"),
                "disposition": "attachment",
            }
            for i, img in enumerate(data["adjuntos"])
        ]

    resp = requests.post(
        "https://api.sendgrid.com/v3/mail/send",
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        json=payload,
        timeout=15,
    )
    if resp.status_code not in (200, 202):
        raise ValueError(f"SendGrid error {resp.status_code}: {resp.text[:200]}")


EMAIL_CONFIRM_KEYBOARD = InlineKeyboardMarkup([[
    InlineKeyboardButton("✅ Enviar",   callback_data="email_confirm"),
    InlineKeyboardButton("✏️ Editar",   callback_data="email_edit"),
    InlineKeyboardButton("❌ Cancelar", callback_data="email_cancel"),
]])

EMAIL_EDIT_KEYBOARD = InlineKeyboardMarkup([[
    InlineKeyboardButton("👤 Destinatario", callback_data="email_edit_para"),
    InlineKeyboardButton("📌 Asunto",       callback_data="email_edit_asunto"),
    InlineKeyboardButton("📝 Cuerpo",       callback_data="email_edit_cuerpo"),
]])


async def artifact_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not context.args:
        await update.message.reply_text(ARTIFACT_HELP, parse_mode="Markdown")
        return

    artifact_type = context.args[0].lower()
    description = " ".join(context.args[1:]) if len(context.args) > 1 else ""

    if artifact_type not in ARTIFACT_PROMPTS:
        await update.message.reply_text(
            f"Tipo `{artifact_type}` no reconocido. Usa: `html`, `excel`, `pdf`, `gantt` o `pptx`.",
            parse_mode="Markdown"
        )
        return

    if not description:
        await update.message.reply_text(
            f"Agrega una descripción. Ejemplo:\n`/artifact {artifact_type} OEE equipos cosecha semana 23`",
            parse_mode="Markdown"
        )
        return

    await update.message.reply_text(f"⏳ Generando *{artifact_type}*...", parse_mode="Markdown")

    description = _build_artifact_description(description, context)
    await _render_artifact(artifact_type, description, _make_reply_fn(update.message), context)


SUPPORTED_DOCS = {".pdf", ".docx", ".xlsx", ".pptx"}

def extract_pptx(data: bytes) -> str:
    """Extrae texto, tablas e imágenes de un PowerPoint."""
    from pptx import Presentation
    from pptx.util import Pt

    prs   = Presentation(io.BytesIO(data))
    parts = []
    total = len(prs.slides)
    parts.append(f"--- DOCUMENTO POWERPOINT: {total} diapositivas en total ---")

    headings = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = []
        title_text  = ""

        # Layout / título de la diapositiva
        if slide.shapes.title and slide.shapes.title.text.strip():
            title_text = slide.shapes.title.text.strip()
            headings.append(title_text)
            slide_parts.append(f"## {title_text}")

        # Recorrer todas las formas
        for shape in slide.shapes:
            # Texto (excluyendo título ya procesado)
            if shape.has_text_frame and shape != slide.shapes.title:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        slide_parts.append(txt)

            # Tablas
            if shape.has_table:
                slide_parts.append(f"[Tabla diapositiva {i}]")
                for row in shape.table.rows:
                    slide_parts.append(" | ".join(
                        cell.text.strip() for cell in row.cells
                    ))

        if slide_parts:
            parts.append(f"\n[Diapositiva {i}/{total}]")
            parts.extend(slide_parts)

    if headings:
        parts.insert(1, "--- SECCIONES: " + " | ".join(headings[:20]) + " ---")

    full = "\n".join(parts)
    parts.insert(1 if not headings else 2,
                 f"--- Total caracteres extraídos: {len(full)} ---")
    return "\n".join(parts)


def extract_pdf(data: bytes) -> str:
    """Extrae texto de un PDF con metadata de estructura."""
    parts = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        total_pages = len(pdf.pages)
        parts.append(f"--- DOCUMENTO PDF: {total_pages} páginas en total ---")
        for i, page in enumerate(pdf.pages, 1):
            page_parts = []
            text = page.extract_text()
            if text:
                page_parts.append(text.strip())
            for table in page.extract_tables():
                for row in table:
                    page_parts.append(" | ".join(str(c or "") for c in row))
            if page_parts:
                parts.append(f"\n[Página {i}/{total_pages}]")
                parts.extend(page_parts)
    full = "\n".join(parts)
    total_chars = len(full)
    parts.insert(1, f"--- Total caracteres extraídos: {total_chars} ---")
    return "\n".join(parts)


def extract_docx(data: bytes) -> str:
    """Extrae texto, estructura y tablas de un documento Word."""
    doc = DocxDocument(io.BytesIO(data))
    parts = []

    # Metadata de estructura
    headings = [p.text.strip() for p in doc.paragraphs if p.style.name.startswith("Heading") and p.text.strip()]
    n_paras  = sum(1 for p in doc.paragraphs if p.text.strip())
    n_tables = len(doc.tables)
    parts.append(f"--- DOCUMENTO WORD: {n_paras} párrafos, {n_tables} tablas ---")
    if headings:
        parts.append("--- SECCIONES: " + " | ".join(headings[:20]) + " ---")

    # Contenido con marcadores de sección
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        if para.style.name.startswith("Heading"):
            parts.append(f"\n## {para.text.strip()}")
        else:
            parts.append(para.text.strip())

    # Tablas con header explícito
    for t_idx, table in enumerate(doc.tables, 1):
        parts.append(f"\n[Tabla {t_idx}]")
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))

    full = "\n".join(parts)
    parts.insert(2, f"--- Total caracteres extraídos: {len(full)} ---")
    return "\n".join(parts)


def extract_xlsx(data: bytes) -> tuple[str, dict]:
    """
    Extrae datos de un Excel.
    Retorna (resumen_texto, datos_estructurados) donde datos_estructurados
    incluye estadísticas y top-N filas listas para usar en HTML/tablas.
    """
    from collections import defaultdict, Counter

    wb  = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    txt = []
    structured = {}   # hoja → {headers, sample, stats, totals}

    for sheet_name in wb.sheetnames[:3]:
        ws = wb[sheet_name]
        headers_row = next(ws.iter_rows(min_row=1, max_row=1, values_only=True))
        headers = [str(h) if h is not None else f"Col{i}" for i, h in enumerate(headers_row)]

        # Lee todas las filas (hasta 5.000 para estadísticas)
        all_rows = []
        for row in ws.iter_rows(min_row=2, max_row=min(ws.max_row, 5001), values_only=True):
            if any(c is not None for c in row):
                all_rows.append(list(row))

        total = len(all_rows)
        txt.append(f"=== Hoja: {sheet_name} ({total} filas, {len(headers)} columnas) ===")
        txt.append(" | ".join(headers))

        # Muestra primeras 25 filas como texto
        for row in all_rows[:25]:
            txt.append(" | ".join(str(c) if c is not None else "" for c in row))
        if total > 25:
            txt.append(f"... ({total - 25} filas adicionales)")

        # Resumen de estadísticas SOBRE TODAS LAS FILAS — incluido en el texto
        # para que Claude lo use en análisis y preguntas de seguimiento
        txt.append("\n--- ESTADÍSTICAS REALES (todas las filas) ---")

        # Estadísticas por columna categórica (máx. 20 valores únicos)
        stats = {}
        for col_i, col_name in enumerate(headers):
            vals = [row[col_i] for row in all_rows if col_i < len(row) and row[col_i] not in (None, "", "<Null>", "None")]
            if not vals:
                continue
            # Numérica
            nums = []
            for v in vals:
                try:
                    nums.append(float(str(v).replace(",", ".")))
                except (ValueError, TypeError):
                    pass
            if len(nums) > len(vals) * 0.5:
                stats[col_name] = {
                    "tipo": "num",
                    "total": len(nums),
                    "suma": round(sum(nums), 2),
                    "min": round(min(nums), 2),
                    "max": round(max(nums), 2),
                    "prom": round(sum(nums) / len(nums), 2),
                }
            elif len(set(str(v) for v in vals)) <= 20:
                cnt = Counter(str(v) for v in vals)
                stats[col_name] = {"tipo": "cat", "frecuencias": dict(cnt.most_common(15))}

        # Muestra representativa: top-20 por la columna numérica de mayor promedio
        # (excluye columnas tipo ID donde max ≈ total de filas)
        num_col = None
        best_score, best_i = -1, None
        for i, h in enumerate(headers):
            if h in stats and stats[h]["tipo"] == "num":
                s = stats[h]
                # Columnas tipo ID tienen prom ≈ total/2 y valores únicos = total
                # Columnas de métricas tienen mayor dispersión relativa
                if s["total"] > 0 and s["max"] > 0:
                    score = s["prom"] / max(s["max"], 1)   # < 1 siempre; IDs tienden a 0.5
                    # Preferir columnas donde la media es > 10% del máximo (métricas reales)
                    if score > best_score and s["max"] < 1e8:
                        best_score, best_i = score, i
        num_col = best_i
        if num_col is not None:
            try:
                top20 = sorted(all_rows, key=lambda r: float(str(r[num_col] or 0).replace(",", ".")), reverse=True)[:20]
            except Exception:
                top20 = all_rows[:20]
        else:
            top20 = all_rows[:20]

        # Vuelca estadísticas en el texto para que Claude las use en análisis
        for col_name, s in stats.items():
            if s["tipo"] == "num":
                txt.append(f"{col_name}: total={s['total']}, suma={s['suma']}, min={s['min']}, max={s['max']}, prom={s['prom']}")
            else:
                freq_str = ", ".join(f"{k}: {v}" for k, v in list(s["frecuencias"].items())[:10])
                txt.append(f"{col_name} ({sum(s['frecuencias'].values())} registros): {freq_str}")

        structured[sheet_name] = {
            "headers":  headers,
            "total_filas": total,
            "muestra_top20": [[str(c) if c is not None else "" for c in r] for r in top20],
            "stats":    stats,
        }

    return "\n".join(txt), structured


async def handle_audio(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Transcribe un mensaje de voz o audio con Groq Whisper y lo procesa con Claude."""
    # Soporta mensajes de voz (micrófono) y archivos de audio
    voice = update.message.voice or update.message.audio
    if not voice:
        return

    await update.message.reply_text("🎙️ Transcribiendo audio...")

    file = await context.bot.get_file(voice.file_id)
    file_bytes = bytes(await file.download_as_bytearray())

    # Groq Whisper requiere un archivo con nombre y extensión
    ext = ".ogg" if update.message.voice else ".mp3"
    try:
        transcription = groq_client.audio.transcriptions.create(
            model="whisper-large-v3",
            file=(f"audio{ext}", io.BytesIO(file_bytes)),
            language="es",
            response_format="text",
        )
        transcript = transcription.strip() if isinstance(transcription, str) else transcription.text.strip()
    except Exception as e:
        await update.message.reply_text(f"⚠️ Error al transcribir el audio: {_safe_err(e)}")
        return

    if not transcript:
        await update.message.reply_text("⚠️ No pude entender el audio. Intenta de nuevo.")
        return

    await update.message.reply_text(f"🗣️ <b>Transcripción:</b> <i>{_html.escape(transcript)}</i>", parse_mode="HTML")

    # Si está en modo notas, guardar la transcripción como nota
    if context.user_data.get("modo_notas"):
        notas = context.user_data.get("notas", [])
        notas.append({"texto": transcript, "fecha": datetime.now().strftime("%d/%m %H:%M"), "n": len(notas) + 1})
        context.user_data["notas"] = notas
        await update.message.reply_text(
            _notas_status_text(notas), parse_mode="Markdown", reply_markup=NOTAS_KEYBOARD
        )
        return

    if context.user_data.get("nlm_mode"):
        context.user_data.pop("nlm_mode")
        await _handle_nlm_query(update, context, transcript)
        return

    await update.message.reply_text("🤖 Analizando con los agentes...")

    try:
        history  = context.user_data.get("history", [])
        try:
            rag_ctx = rag.build_context(transcript) if rag.col and rag.col.count() > 0 else ""
        except Exception:
            rag_ctx = ""
        system   = SYSTEM_PROMPT + rag_ctx
        reply = claude_response(system, transcript, max_tokens=600,
                                model=get_model_for_msg(context, transcript), history=history)
        push_history(context, transcript, reply)
        context.user_data["last_analysis"] = reply

        try:
            await update.message.reply_text(
                fmt(reply) + "\n\n🎨 <i>¿Generar un artefacto visual con esto?</i>",
                reply_markup=ARTIFACT_KEYBOARD,
                parse_mode="HTML"
            )
        except Exception:
            # Fallback sin Markdown si la respuesta tiene caracteres problemáticos
            await update.message.reply_text(
                reply + "\n\n🎨 ¿Generar un artefacto visual con esto?",
                reply_markup=ARTIFACT_KEYBOARD
            )
    except Exception as e:
        await update.message.reply_text(f"⚠️ Error al procesar: {_safe_err(e)}")


_MAX_DOC_BYTES = 20 * 1024 * 1024  # 20 MB

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    doc = update.message.document
    filename = (doc.file_name or "").lower()
    ext = next((e for e in SUPPORTED_DOCS if filename.endswith(e)), None)

    if not ext:
        await update.message.reply_text(
            "📎 Solo proceso *PDF*, *Word (.docx)* y *Excel (.xlsx)*.",
            parse_mode="Markdown"
        )
        return

    if doc.file_size and doc.file_size > _MAX_DOC_BYTES:
        await update.message.reply_text(
            f"⚠️ El archivo supera el límite de 20 MB ({doc.file_size // (1024*1024)} MB)."
        )
        return

    await update.message.reply_text("📂 Leyendo archivo...")

    file = await context.bot.get_file(doc.file_id)
    file_bytes = bytes(await file.download_as_bytearray())

    try:
        if ext == ".pdf":
            content = extract_pdf(file_bytes)
            structured_data = {}
            tipo = "PDF"
        elif ext == ".docx":
            content = extract_docx(file_bytes)
            structured_data = {}
            tipo = "Word"
        elif ext == ".pptx":
            content = extract_pptx(file_bytes)
            structured_data = {}
            tipo = "PowerPoint"
        else:
            content, structured_data = extract_xlsx(file_bytes)
            tipo = "Excel"
    except Exception as e:
        await update.message.reply_text(f"⚠️ No pude leer el archivo: {_safe_err(e)}")
        return

    if not content.strip():
        await update.message.reply_text("⚠️ El archivo no tiene contenido extraíble.")
        return

    await update.message.reply_text("🤖 Analizando con los agentes...")

    # 10.000 chars para todos los tipos — Excel incluye stats, PDF/Word incluyen estructura
    prompt = (
        f"Analiza este documento {tipo} en el contexto operacional forestal de Arauco. "
        f"Identifica datos clave, KPIs, procesos, problemas u oportunidades de mejora.\n\n"
        f"{content[:10000]}"
    )
    history  = context.user_data.get("history", [])
    analysis = claude_response(SYSTEM_PROMPT, prompt, max_tokens=800, model=get_model(context), history=history)
    push_history(context, f"[Análisis de {tipo}: {doc.file_name}]", analysis)
    context.user_data["last_analysis"] = analysis
    # Guarda datos estructurados (Excel) y contenido raw (PDF/Word/Excel)
    # para que el artifact HTML pueda generar tablas con datos reales
    context.user_data["structured_data"] = structured_data   # dict (Excel) o {} (PDF/Word)
    context.user_data["doc_content"]     = content[:10000]   # texto raw del documento
    context.user_data["doc_tipo"]        = tipo               # "PDF" | "Word" | "Excel"
    # Guarda contenido completo para indexar si el usuario lo solicita
    context.user_data["pending_index"] = {"filename": doc.file_name, "content": content}

    header = f"📄 <b>{doc.file_name}</b>\n\n"
    try:
        await update.message.reply_text(
            header + fmt(analysis) + "\n\n🎨 <i>¿Qué deseas hacer con este documento?</i>",
            reply_markup=ARTIFACT_KEYBOARD,
            parse_mode="HTML"
        )
    except Exception:
        await send_reply(update, f"{doc.file_name}\n\n{analysis}", reply_markup=ARTIFACT_KEYBOARD, context=context)


async def email_confirm_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    await query.edit_message_reply_markup(reply_markup=None)

    if query.data == "email_cancel":
        await query.message.reply_text("❌ Envío cancelado.")
        context.user_data.pop("pending_email", None)
        return

    if query.data == "email_edit":
        await query.message.reply_text("¿Qué quieres editar?", reply_markup=EMAIL_EDIT_KEYBOARD)
        return

    if query.data.startswith("email_edit_"):
        campo = query.data.replace("email_edit_", "")
        labels = {"para": "destinatario", "asunto": "asunto", "cuerpo": "cuerpo del correo"}
        if campo not in labels:
            await query.answer("Campo no válido.", show_alert=True)
            return
        context.user_data["editing_email_field"] = campo
        await query.message.reply_text(f"Escribe el nuevo {labels[campo]}:")
        return

    data = context.user_data.get("pending_email")
    if not data:
        await query.message.reply_text("⚠️ No hay correo pendiente.")
        return

    data["adjuntos"] = context.user_data.pop("email_attachments", [])
    await query.message.reply_text("⏳ Enviando correo...")
    try:
        send_email_outlook(data)
        context.user_data.pop("pending_email", None)
        context.user_data.pop("email_n_adjuntos", None)
        n_adj = len(data["adjuntos"])
        adj_txt = f"\n📎 {n_adj} imagen(es) adjunta(s)" if n_adj else ""
        await query.message.reply_text(
            f"✅ Correo enviado a <b>{_html.escape(data['para'])}</b>\n"
            f"<b>Asunto:</b> {_html.escape(data['asunto'])}{adj_txt}",
            parse_mode="HTML"
        )
    except Exception as e:
        await query.message.reply_text(f"⚠️ Error al enviar: {_safe_err(e)}")


async def rag_index_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Indexa el documento actual en RAG con respuesta compacta."""
    query = update.callback_query
    await query.answer()

    pending = context.user_data.get("pending_index")
    if not pending:
        await query.message.reply_text("⚠️ No hay documento pendiente de indexar.")
        return

    await query.edit_message_reply_markup(reply_markup=None)

    try:
        n = rag.index_document(pending["content"], pending["filename"])
        context.user_data.pop("pending_index", None)
        docs = rag.list_documents()
        await query.message.reply_text(
            f"✅ <b>{_html.escape(pending['filename'])}</b> indexado ({n} fragmentos)\n"
            f"Base: {len(docs)} documento(s). Usa /notebookrag para consultar.",
            parse_mode="HTML"
        )
    except Exception as e:
        await query.message.reply_text(f"⚠️ Error al indexar: {_safe_err(e)}")


async def documentos_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Lista documentos indexados y activa NotebookRAG para consulta inmediata."""
    try:
        docs  = rag.list_documents()
        total = rag.col.count() if rag.col else 0
        if not docs:
            await update.message.reply_text(
                "📚 Base de conocimiento vacía.\n"
                "Sube un PDF, Word o Excel y presiona <b>📥 RAG</b>.",
                parse_mode="HTML"
            )
            return
        doc_list = "\n".join(f"• {d}" for d in docs)
        context.user_data["nlm_mode"] = True
        await update.message.reply_text(
            f"📖 <b>NotebookRAG</b> — {len(docs)} documento(s) · {total} fragmentos\n\n"
            f"{doc_list}\n\n"
            "Escribe tu pregunta:",
            parse_mode="HTML"
        )
    except Exception as e:
        await update.message.reply_text(f"⚠️ Error: {_safe_err(e)}")


async def buscar_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Prueba la búsqueda RAG directamente sin pasar por Claude."""
    query_text = " ".join(context.args) if context.args else ""
    if not query_text:
        await update.message.reply_text(
            "Uso: <code>/buscar texto a buscar</code>\n"
            "Ejemplo: <code>/buscar rangos de pendiente terreno</code>",
            parse_mode="HTML"
        )
        return
    try:
        chunks = rag.query(query_text)
        if not chunks:
            await update.message.reply_text(
                f"🔍 Sin resultados para: <i>{_html.escape(query_text)}</i>\n\n"
                "Verifica con /documentos que el documento esté indexado.",
                parse_mode="HTML"
            )
            return
        resp = f"🔍 <b>Resultados para:</b> <i>{_html.escape(query_text)}</i>\n\n"
        for i, c in enumerate(chunks, 1):
            resp += (f"<b>{i}. {_html.escape(c['filename'])}</b> "
                     f"(relevancia: {c['score']})\n"
                     f"<i>{_html.escape(c['text'][:300])}...</i>\n\n")
        await update.message.reply_text(resp, parse_mode="HTML")
    except Exception as e:
        await update.message.reply_text(f"⚠️ Error: {_safe_err(e)}")


async def indexar_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Instrucción para indexar documentos."""
    await update.message.reply_text(
        "📚 *Indexar documento*\n\n"
        "Envía un archivo *PDF*, *Word (.docx)* o *Excel (.xlsx)* "
        "y al finalizar el análisis presiona el botón *📚 Indexar en RAG*.\n\n"
        "El documento quedará disponible para que los agentes lo consulten "
        "automáticamente al responder preguntas.",
        parse_mode="Markdown"
    )


from telegram import BotCommand

START_TEXT = """🌲 *Arauco — Subgerencia de Mejora Continua*

Soy el asistente digital de tu equipo. Integro tres agentes especializados coordinados por el Subgerente de Mejora Continua:

🏭 *Agente EO — Excelencia Operacional*
Lean, GEMBA, KAIZEN, BPMN 2.0, KPIs, OEE, A3/PDCA, rediseño de procesos forestales

🤖 *Agente IA — Inteligencia Artificial*
Modelos predictivos (ML/XGBoost), GenAI con Claude API, LangGraph, cartografía con IA, dashboards HTML

⚙️ *Agente TD — Transformación Digital*
Integraciones SAP/SGL/Planex/Forest Data, telemetría de maquinaria forestal (Tigercat, John Deere, Develon), ETL, arquitecturas de datos

🧭 *Orquestador (Subgerente MC)*
Coordina los agentes, sintetiza resultados y entrega análisis ejecutivos estilo McKinsey/BCG

---

*Comandos disponibles:*
/spec — Especificación de iniciativa
/plan — Plan de ejecución
/build — Construcción de solución
/test — Validación operacional
/review — Revisión crítica
/ship — Lanzamiento a operación
/artifact — Genera Excel, gráfico o dashboard HTML

También puedes enviar una imagen, PDF, Word o Excel y los agentes lo analizarán."""

async def start_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(START_TEXT, parse_mode="Markdown")


def _track(context, message_id: int):
    """Registra un message_id enviado por el bot para poder borrarlo en reset."""
    if context is None:
        return
    sent = context.user_data.setdefault("sent_messages", [])
    sent.append(message_id)
    if len(sent) > 500:
        context.user_data["sent_messages"] = sent[-500:]


async def _do_reset(chat_id: int, context) -> None:
    """Borra mensajes del bot (best-effort ≤48h), limpia estado, envía divisor."""
    for mid in context.user_data.get("sent_messages", []):
        try:
            await context.bot.delete_message(chat_id, mid)
        except Exception:
            pass

    model = context.user_data.get("model", DEFAULT_MODEL)
    context.user_data.clear()
    context.user_data["model"] = model
    context.user_data["sent_messages"] = []

    now = datetime.now().strftime("%d %b · %H:%M")
    await context.bot.send_message(
        chat_id,
        f"━━━━━━━━━━━━━━━━━━━━━━━━━\n🆕  Nueva conversación  ·  {now}\n━━━━━━━━━━━━━━━━━━━━━━━━━"
    )


async def reset_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Muestra confirmación antes de limpiar la conversación."""
    kb = InlineKeyboardMarkup([[
        InlineKeyboardButton("✅ Confirmar", callback_data="reset_confirm"),
        InlineKeyboardButton("❌ Cancelar",  callback_data="reset_cancel"),
    ]])
    msg = await update.message.reply_text(
        "¿Iniciar nueva conversación?\nSe borrará el historial y los documentos cargados.",
        reply_markup=kb
    )
    _track(context, msg.message_id)


async def reset_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Maneja confirmación/cancelación del reset."""
    query = update.callback_query
    await query.answer()
    if query.data == "reset_confirm":
        await _do_reset(query.message.chat_id, context)
    else:
        await query.message.delete()


async def handle_location(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Guarda una ubicación GPS como nota o la adjunta a la última nota."""
    loc  = update.message.location
    lat  = loc.latitude
    lon  = loc.longitude
    maps = f"https://maps.google.com/?q={lat},{lon}"
    gps_line = f"📍 {lat:.6f}, {lon:.6f} — {maps}"

    # Caso 1: adjuntar a la última nota (botón "📍 + ubicación a esta nota")
    target_idx = context.user_data.pop("agregar_ubicacion_a_nota", None)
    if target_idx is not None:
        notas = context.user_data.get("notas", [])
        if notas and target_idx <= len(notas):
            notas[target_idx - 1]["texto"] += f"\n{gps_line}"
            context.user_data["notas"] = notas
            nota = notas[target_idx - 1]
            preview = nota["texto"][:80] + ("…" if len(nota["texto"]) > 80 else "")
            await update.message.reply_text(
                f"📍 *Ubicación añadida a Nota {target_idx}*\n_{preview}_",
                parse_mode="Markdown", reply_markup=NOTAS_KEYBOARD
            )
        else:
            await update.message.reply_text("⚠️ No se encontró la nota.", reply_markup=NOTAS_KEYBOARD)
        return

    # Caso 2: modo notas activo → nueva nota con solo la ubicación
    if context.user_data.get("modo_notas"):
        notas = context.user_data.get("notas", [])
        notas.append({"texto": gps_line, "fecha": datetime.now().strftime("%d/%m %H:%M"), "n": len(notas) + 1})
        context.user_data["notas"] = notas
        await update.message.reply_text(
            _notas_status_text(notas), parse_mode="Markdown", reply_markup=NOTAS_KEYBOARD
        )
        return

    # Caso 3: fuera de modo notas
    await update.message.reply_text(
        f"📍 *Ubicación recibida*\n`{lat:.6f}, {lon:.6f}`\n[Ver en Google Maps]({maps})\n\n"
        "_Activa el modo notas para guardar ubicaciones automáticamente._",
        parse_mode="Markdown"
    )


async def post_init(application):
    await application.bot.set_my_commands([
        BotCommand("start",       "🌲 Qué soy y cómo funciono"),
        BotCommand("reset",       "🔄 Reiniciar conversación"),
        BotCommand("modelo",      "🤖 Seleccionar modelo LLM"),
        BotCommand("indexar",     "📚 Cómo indexar documentos en RAG"),
        BotCommand("notebookrag", "📖 Consultar base de conocimiento Arauco"),
        BotCommand("spec",        "📋 Especificación de iniciativa forestal"),
        BotCommand("plan",        "🗺️ Plan de ejecución forestal"),
        BotCommand("build",       "🔨 Construcción de solución forestal"),
        BotCommand("test",        "🧪 Validación en operación forestal"),
        BotCommand("review",      "🔍 Revisión crítica forestal"),
        BotCommand("ship",        "🚀 Lanzamiento a operación forestal"),
        BotCommand("automation",  "⚙️ Automatización de procesos forestales"),
        BotCommand("connectivity","📡 Conectividad en predios remotos"),
        BotCommand("facilitation","🤝 Facilitación de talleres Lean"),
        BotCommand("telemetry",   "📊 Telemetría de maquinaria forestal"),
        BotCommand("artifact",    "🎨 Genera HTML, Excel o gráfico PNG"),
    ])

app = (
    ApplicationBuilder()
    .token(os.environ["TELEGRAM_TOKEN"])
    .persistence(PicklePersistence(filepath="/tmp/bot_persistence"))
    .post_init(post_init)
    .build()
)

app.add_handler(CommandHandler("start",      start_handler))
app.add_handler(CommandHandler("reset",      reset_handler))
app.add_handler(CallbackQueryHandler(reset_callback, pattern="^reset_"))
app.add_handler(CommandHandler("modelo",     modelo_handler))
app.add_handler(CommandHandler("indexar",    indexar_handler))
app.add_handler(CommandHandler("notebookrag", documentos_handler))
app.add_handler(CommandHandler("documentos",  documentos_handler))
app.add_handler(CommandHandler("buscar",     buscar_handler))
app.add_handler(CallbackQueryHandler(modelo_callback,    pattern="^mdl_"))
app.add_handler(CallbackQueryHandler(rag_index_callback,   pattern="^rag_index$"))

app.add_handler(CallbackQueryHandler(email_confirm_callback, pattern="^email_(confirm|cancel|edit.*)$"))
app.add_handler(CallbackQueryHandler(notas_callback,        pattern="^notas_"))
app.add_handler(CallbackQueryHandler(image_callback,        pattern="^img_"))

for skill in SKILL_PROMPTS:
    app.add_handler(CommandHandler(skill, skill_handler))

app.add_handler(CommandHandler("artifact", artifact_handler))
app.add_handler(CallbackQueryHandler(artifact_callback, pattern="^art_"))
app.add_handler(MessageHandler(filters.LOCATION, handle_location))
app.add_handler(MessageHandler(filters.VOICE, handle_audio))
app.add_handler(MessageHandler(filters.AUDIO, handle_audio))
app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))
app.add_handler(MessageHandler(filters.PHOTO, handle_image))


app.run_polling(drop_pending_updates=True)
